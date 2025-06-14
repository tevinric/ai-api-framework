from flask import jsonify, request, g, make_response
from apis.utils.tokenService import TokenService
from apis.utils.databaseService import DatabaseService
from apis.utils.logMiddleware import api_logger
from apis.utils.balanceMiddleware import check_balance
from apis.utils.config import get_azure_blob_client, ensure_container_exists
from apis.utils.fileService import FileService
from apis.utils.llmServices import gpt4o_mini_service, gpt4o_service
import logging
import uuid
import pytz
from datetime import datetime
import os
import io
import json
import base64
import time
import re
import requests

# For document processing
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation
import docx
import pandas as pd
import tempfile

# CONFIGURE LOGGING
logger = logging.getLogger(__name__)

# Define container for file uploads
FILE_UPLOAD_CONTAINER = os.environ.get("AZURE_STORAGE_UPLOAD_CONTAINER", "file-uploads")
STORAGE_ACCOUNT = os.environ.get("AZURE_STORAGE_ACCOUNT")
BASE_BLOB_URL = f"https://{STORAGE_ACCOUNT}.blob.core.windows.net/{FILE_UPLOAD_CONTAINER}"

from apis.utils.config import create_api_response

def extract_text_from_pdf(file_path):
    """Extract text from PDF files using PyMuPDF (faster and better for large files)"""
    try:
        logger.info(f"Opening PDF file: {file_path}")
        doc = fitz.open(file_path)
        total_pages = len(doc)
        logger.info(f"PDF has {total_pages} pages")
        
        if total_pages == 0:
            raise ValueError("PDF file has no pages")
            
        text_content = []
        
        # Process each page without truncation
        for page_num in range(total_pages):
            page = doc.load_page(page_num)
            page_text = page.get_text()
            text_content.append(page_text)
            
            # Log first few characters of first and last page for debugging
            if page_num == 0 or page_num == total_pages - 1:
                preview = page_text[:100].replace('\n', ' ').strip()
                logger.info(f"Page {page_num+1} preview: {preview}...")
        
        full_text = "\n".join(text_content)
        logger.info(f"Extracted {len(full_text)} characters of text from PDF")
        
        return full_text, total_pages
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        raise

def extract_text_from_docx(file_path):
    """Extract text from DOCX files"""
    try:
        logger.info(f"Opening DOCX file: {file_path}")
        doc = docx.Document(file_path)
        paragraphs = []
        
        # Extract paragraphs
        logger.info(f"DOCX has {len(doc.paragraphs)} paragraphs")
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
                
        # Extract text from tables
        logger.info(f"DOCX has {len(doc.tables)} tables")
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text)
        
        # Number of paragraphs as a proxy for "pages"
        # This is not perfect but gives an estimate
        total_pages = max(1, len(paragraphs) // 4)  # Rough estimate: ~4 paragraphs per page
        
        full_text = "\n".join(paragraphs)
        logger.info(f"Extracted {len(full_text)} characters from DOCX")
        
        # Log a preview of the text for debugging
        if paragraphs:
            preview = paragraphs[0][:100].replace('\n', ' ').strip()
            logger.info(f"Text preview: {preview}...")
            
        return full_text, total_pages
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise

def extract_text_from_pptx(file_path):
    """Extract text from PPTX files"""
    try:
        logger.info(f"Opening PPTX file: {file_path}")
        presentation = Presentation(file_path)
        total_slides = len(presentation.slides)
        logger.info(f"PPTX has {total_slides} slides")
        
        text_content = []
        
        # Process each slide
        for slide_idx, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            slide_content = "\n".join(slide_text)
            text_content.append(f"--- Slide {slide_idx + 1} ---\n{slide_content}")
            
            # Log preview of first and last slide for debugging
            if slide_idx == 0 or slide_idx == total_slides - 1:
                preview = slide_content[:100].replace('\n', ' ').strip()
                logger.info(f"Slide {slide_idx+1} preview: {preview}...")
        
        full_text = "\n\n".join(text_content)
        logger.info(f"Extracted {len(full_text)} characters from PPTX")
        
        return full_text, total_slides
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise

def extract_text_from_xlsx(file_path):
    """Extract text from Excel files"""
    try:
        logger.info(f"Opening Excel file: {file_path}")
        
        # Read all sheets
        all_sheets = pd.read_excel(file_path, sheet_name=None)
        total_sheets = len(all_sheets)
        logger.info(f"Excel file has {total_sheets} sheets")
        
        text_content = []
        
        # Process each sheet
        for sheet_name, df in all_sheets.items():
            logger.info(f"Processing sheet '{sheet_name}' with {len(df)} rows and {len(df.columns)} columns")
            
            # Add a clear sheet header
            text_content.append(f"=== Sheet: {sheet_name} ===")
            
            # Convert dataframe to string representation
            sheet_text = df.to_string(index=False)
            text_content.append(sheet_text)
            text_content.append("\n")
            
            # Log preview for debugging
            if not df.empty:
                preview = str(df.head(1)).replace('\n', ' ')[:100]
                logger.info(f"Sheet '{sheet_name}' preview: {preview}...")
        
        full_text = "\n".join(text_content)
        logger.info(f"Extracted {len(full_text)} characters from Excel file")
        
        return full_text, total_sheets
    except Exception as e:
        logger.error(f"Error extracting text from Excel: {str(e)}")
        raise

def extract_token_info_from_llm_response(response_data):
    """
    Directly extract token usage information from LLM responses
    This handles various formats that different LLM providers might use
    """
    token_info = {
        "prompt_tokens": 0,
        "completion_tokens": 0,
        "total_tokens": 0
    }
    
    if not response_data:
        return token_info
    
    # If it's already a dictionary, try to extract token fields
    if isinstance(response_data, dict):
        # Try standard token fields
        if "usage" in response_data and isinstance(response_data["usage"], dict):
            usage = response_data["usage"]
            if "prompt_tokens" in usage:
                token_info["prompt_tokens"] = usage["prompt_tokens"]
            if "completion_tokens" in usage:
                token_info["completion_tokens"] = usage["completion_tokens"]
            if "total_tokens" in usage:
                token_info["total_tokens"] = usage["total_tokens"]
            return token_info
        
        # Try direct token fields
        token_field_mapping = {
            "input_tokens": "prompt_tokens",
            "prompt_tokens": "prompt_tokens",
            "completion_tokens": "completion_tokens",
            "output_tokens": "total_tokens",
            "total_tokens": "total_tokens"
        }
        
        for source_field, target_field in token_field_mapping.items():
            if source_field in response_data and isinstance(response_data[source_field], (int, float)):
                token_info[target_field] = response_data[source_field]
        
        # If we have prompt and completion but no total, calculate it
        if token_info["prompt_tokens"] > 0 or token_info["completion_tokens"] > 0:
            if token_info["total_tokens"] == 0:
                token_info["total_tokens"] = token_info["prompt_tokens"] + token_info["completion_tokens"]
        
        return token_info
    
    # If it's a string that looks like JSON, try to parse it
    if isinstance(response_data, str) and response_data.strip().startswith('{') and response_data.strip().endswith('}'):
        try:
            data = json.loads(response_data)
            return extract_token_info_from_llm_response(data)
        except json.JSONDecodeError:
            return token_info
    
    # If we get here, we couldn't extract token info
    return token_info

def parse_llm_response_json(content):
    """
    Parse potential JSON in LLM responses and extract the actual content
    This handles cases where the LLM returns a JSON object with the actual content in a field
    """
    # If it's not a string, convert it to a string
    if not isinstance(content, str):
        return str(content), {}
    
    # If it doesn't look like JSON, return as is
    if not (content.strip().startswith('{') and content.strip().endswith('}')):
        return content, {}
    
    try:
        # Try to parse as JSON
        json_obj = json.loads(content)
        
        # If it's not a dict, return as is
        if not isinstance(json_obj, dict):
            return content, {}
        
        # Check for common fields in LLM responses
        extracted_content = None
        metadata = {}
        
        # Look for the actual content
        for field in ["message", "summary", "content", "text", "result"]:
            if field in json_obj and isinstance(json_obj[field], str):
                extracted_content = json_obj[field]
                break
        
        # If we didn't find any of those fields, check if there's only one string field
        if extracted_content is None:
            string_fields = {k: v for k, v in json_obj.items() if isinstance(v, str)}
            if len(string_fields) == 1:
                extracted_content = next(iter(string_fields.values()))
        
        # Extract token metadata if present
        token_fields = {
            "input_tokens": "prompt_tokens",
            "prompt_tokens": "prompt_tokens",
            "completion_tokens": "completion_tokens",
            "output_tokens": "total_tokens",
            "total_tokens": "total_tokens"
        }
        
        for json_field, meta_field in token_fields.items():
            if json_field in json_obj and isinstance(json_obj[json_field], (int, float)):
                metadata[meta_field] = json_obj[json_field]
        
        # If we found content, return it along with metadata
        if extracted_content is not None:
            return extracted_content, metadata
        
        # If we get here, we couldn't extract content, return the original
        return content, {}
    except json.JSONDecodeError:
        # Not valid JSON, return as is
        return content, {}

def extract_summary_from_llm_response(response_data):
    """
    Extract actual summary text from LLM response JSON
    This handles cases where the summary contains a full JSON response
    """
    if not response_data:
        return "No summary available", [], {}
    
    # If response_data is already a string, return it
    if isinstance(response_data, str):
        # Check if it's a JSON string
        parsed_content, metadata = parse_llm_response_json(response_data)
        if parsed_content != response_data:
            # We were able to extract something
            key_points = extract_key_points_from_content(parsed_content)
            return parsed_content, key_points, metadata
        return response_data, [], {}
    
    # Check if we received a JSON object with a 'message' field
    if isinstance(response_data, dict):
        # If it has a message field, that's the actual summary
        if "message" in response_data:
            summary_text = response_data["message"]
            
            # Extract token information
            token_info = {
                "prompt_tokens": response_data.get("input_tokens", 0),
                "completion_tokens": response_data.get("completion_tokens", 0),
                "total_tokens": response_data.get("output_tokens", 0)
            }
            
            # Try to extract key points from the summary
            key_points = extract_key_points_from_content(summary_text)
            
            return summary_text, key_points, token_info
        
        # If it has a summary field, use that
        if "summary" in response_data:
            summary_text = response_data["summary"]
            key_points = response_data.get("key_points", [])
            token_info = response_data.get("tokens", {})
            
            # If summary is still a dict, try to extract it
            if isinstance(summary_text, dict) and "message" in summary_text:
                summary_text = summary_text["message"]
                
            return summary_text, key_points, token_info
    
    # If we get here, we couldn't extract a proper summary
    # Just convert the whole thing to a string
    return str(response_data), [], {}

def extract_key_points_from_content(content, min_points=5):
    """
    Extract key points from content if the LLM didn't properly format them
    This is a fallback mechanism to ensure we always have key points
    """
    logger.info(f"Attempting to extract key points from content of length {len(content)}")
    
    # If content is already a list of key points, return it
    if isinstance(content, list):
        logger.info(f"Content is already a list with {len(content)} items")
        return content
    
    # If content is not a string, convert it
    if not isinstance(content, str):
        logger.warning(f"Content is not a string (type: {type(content)}), converting")
        content = str(content)
    
    extracted_points = []
    
    # Look for bullet points or numbered lists
    bullet_patterns = [
        r'•\s*(.*?)(?=(?:•|\n\n|\Z))',  # Bullet points
        r'^\s*\*\s*(.*?)$',              # Asterisk bullets
        r'^\s*-\s*(.*?)$',               # Dash bullets
        r'^\s*\d+\.\s*(.*?)$',           # Numbered points
        r'Key Point\s*\d*\s*:\s*(.*?)$'  # "Key Point X:" format
    ]
    
    for pattern in bullet_patterns:
        matches = re.findall(pattern, content, re.MULTILINE)
        if matches:
            logger.info(f"Found {len(matches)} key points using pattern: {pattern}")
            for match in matches:
                if match.strip() and len(match.strip()) > 10:  # Minimum length to be considered a key point
                    extracted_points.append(match.strip())
    
    # If we still don't have enough points, use sentences
    if len(extracted_points) < min_points:
        logger.info(f"Only found {len(extracted_points)} key points, extracting from sentences")
        
        # Split into sentences and use the most informative ones as key points
        sentences = re.split(r'(?<=[.!?])\s+', content)
        
        # Filter for sentences that seem informative
        informative_sentences = []
        for sentence in sentences:
            # Clean the sentence
            clean_sentence = sentence.strip()
            
            # Only use sentences that are substantial but not too long
            if 20 <= len(clean_sentence) <= 150 and clean_sentence.endswith(('.', '!', '?')):
                informative_sentences.append(clean_sentence)
        
        # Add the most informative sentences to our key points
        if informative_sentences:
            # Sort by length (longer sentences often have more information)
            informative_sentences.sort(key=len, reverse=True)
            
            # Add sentences until we have enough key points
            for sentence in informative_sentences:
                if sentence not in extracted_points:  # Avoid duplicates
                    extracted_points.append(sentence)
                    if len(extracted_points) >= min_points:
                        break
    
    logger.info(f"Extracted {len(extracted_points)} key points")
    
    # If we still don't have any key points, create generic ones
    if not extracted_points:
        logger.warning("Could not extract any key points, using generic ones")
        extracted_points = ["No specific key points could be identified in the document"]
    
    return extracted_points

def chunk_text(text, max_chunk_size=12000, overlap=500):  # Increased overlap for better context
    """
    Split text into chunks of maximum size with overlap
    Args:
        text: The text to split
        max_chunk_size: Maximum size of each chunk
        overlap: Overlap between chunks
    Returns:
        List of text chunks
    """
    logger.info(f"Chunking text of length {len(text)} with max_chunk_size={max_chunk_size}, overlap={overlap}")
    
    if len(text) <= max_chunk_size:
        logger.info("Text fits in a single chunk")
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + max_chunk_size
        if end < len(text):
            # Find a good breaking point for natural text flow
            # Look for a paragraph break first
            paragraph_break = text.rfind('\n\n', start + max_chunk_size - overlap, end)
            if paragraph_break != -1:
                end = paragraph_break + 2
                logger.info(f"Found paragraph break at position {paragraph_break}")
            else:
                # Look for a sentence break
                sentence_break = text.rfind('. ', start + max_chunk_size - overlap, end)
                if sentence_break != -1:
                    end = sentence_break + 2
                    logger.info(f"Found sentence break at position {sentence_break}")
                else:
                    # If all else fails, look for any space
                    space = text.rfind(' ', start + max_chunk_size - overlap, end)
                    if space != -1:
                        end = space + 1
                        logger.info(f"Found space break at position {space}")
                    else:
                        logger.info(f"No natural break found, using hard cutoff at position {end}")
        
        # Ensure we're not creating tiny chunks at the end
        if end >= len(text) - (overlap / 2) and end < len(text):
            end = len(text)
            logger.info(f"Extending final chunk to include remaining {len(text) - end} characters")
        
        chunk = text[start:end]
        chunks.append(chunk)
        logger.info(f"Created chunk {len(chunks)} with length {len(chunk)}")
        
        # Check for potential data loss
        if end < len(text) and len(chunk) < (max_chunk_size * 0.75):
            logger.warning(f"Chunk {len(chunks)} is significantly smaller than max size, possible content splitting issue")
        
        start = end - overlap
    
    # Verify that all text has been chunked properly
    total_unique_chars = len(text)
    chunked_chars = sum(len(c) for c in chunks) - (overlap * (len(chunks) - 1))
    if chunked_chars < total_unique_chars:
        logger.warning(f"Potential data loss: Original text has {total_unique_chars} chars, chunks cover {chunked_chars} chars")
    
    logger.info(f"Text split into {len(chunks)} chunks")
    return chunks

def format_docx_summary(summary_data):
    """Format summary data for docx with improved JSON handling"""
    # Extract summary, key points and token info from potential JSON response
    summary_text, key_points, token_info = extract_summary_from_llm_response(summary_data)
    
    # Update token info if available in original data
    if "tokens" in summary_data and isinstance(summary_data["tokens"], dict):
        token_info = summary_data["tokens"]
    
    # Use existing key points if available and not empty
    if "key_points" in summary_data and summary_data["key_points"]:
        key_points = summary_data["key_points"]
    
    # If we still don't have key points, extract them from the summary
    if not key_points:
        key_points = extract_key_points_from_content(summary_text)
    
    # Create formatted response with correct structure
    formatted = {
        "summary": summary_text,
        "key_points": key_points,
        "document_structure": summary_data.get("document_structure", {}),
        "pages_processed": summary_data.get("pages_processed", 0),
        "tokens": token_info,
        "model": summary_data.get("model", "")
    }
    
    return formatted

def format_pdf_summary(summary_data):
    """Format summary data for PDF with improved JSON handling"""
    # Extract summary, key points and token info from potential JSON response
    summary_text, key_points, token_info = extract_summary_from_llm_response(summary_data)
    
    # Update token info if available in original data
    if "tokens" in summary_data and isinstance(summary_data["tokens"], dict):
        token_info = summary_data["tokens"]
    
    # Use existing key points if available and not empty
    if "key_points" in summary_data and summary_data["key_points"]:
        key_points = summary_data["key_points"]
    
    # If we still don't have key points, extract them from the summary
    if not key_points:
        key_points = extract_key_points_from_content(summary_text)
    
    # Create formatted response with correct structure
    formatted = {
        "summary": summary_text,
        "key_points": key_points,
        "document_structure": summary_data.get("document_structure", {}),
        "pages_processed": summary_data.get("pages_processed", 0),
        "tokens": token_info,
        "model": summary_data.get("model", "")
    }
    
    return formatted

def format_pptx_summary(summary_data):
    """Format summary data for PowerPoint with improved JSON handling"""
    # Extract summary, key points and token info from potential JSON response
    summary_text, key_points, token_info = extract_summary_from_llm_response(summary_data)
    
    # Update token info if available in original data
    if "tokens" in summary_data and isinstance(summary_data["tokens"], dict):
        token_info = summary_data["tokens"]
    
    # Use existing key points if available and not empty
    if "key_points" in summary_data and summary_data["key_points"]:
        key_points = summary_data["key_points"]
    
    # If we still don't have key points, extract them from the summary
    if not key_points:
        key_points = extract_key_points_from_content(summary_text)
    
    # Create formatted response with correct structure
    formatted = {
        "summary": summary_text,
        "key_points": key_points,
        "slides_processed": summary_data.get("pages_processed", 0),
        "presentation_structure": summary_data.get("document_structure", {}),
        "tokens": token_info,
        "model": summary_data.get("model", "")
    }
    
    return formatted

def format_xlsx_summary(summary_data):
    """Format summary data for Excel with improved JSON handling"""
    # Extract summary, key points and token info from potential JSON response
    summary_text, key_points, token_info = extract_summary_from_llm_response(summary_data)
    
    # Update token info if available in original data
    if "tokens" in summary_data and isinstance(summary_data["tokens"], dict):
        token_info = summary_data["tokens"]
    
    # Use existing key points if available and not empty
    if "key_points" in summary_data and summary_data["key_points"]:
        key_points = summary_data["key_points"]
    
    # If we still don't have key points, extract them from the summary
    if not key_points:
        key_points = extract_key_points_from_content(summary_text)
    
    # Create formatted response with correct structure
    formatted = {
        "summary": summary_text,
        "key_points": key_points,
        "sheets_processed": summary_data.get("pages_processed", 0),
        "data_insights": summary_data.get("document_structure", {}),
        "tokens": token_info,
        "model": summary_data.get("model", "")
    }
    
    return formatted

def chunk_and_summarize(text, total_pages, summary_options):
    """
    Process text in chunks and combine summaries using the llmServices module
    Args:
        text: Full text content
        total_pages: Number of pages in document
        summary_options: Configuration for summarization
    Returns:
        Combined summary and token usage
    """
    import time  # Add time import for sleep functionality
    
    # Use a smaller default chunk size to prevent timeouts in production
    chunk_size = min(summary_options.get('chunk_size', 8000), 8000)  # Cap at 8000 to prevent timeouts
    logger.info(f"Using chunk size: {chunk_size}")
    
    chunks = chunk_text(text, max_chunk_size=chunk_size)
    logger.info(f"Document split into {len(chunks)} chunks for processing")
    
    all_summaries = []
    token_usage = {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0}
    models_used = []  # Track models used for each chunk
    
    for idx, chunk in enumerate(chunks):
        logger.info(f"Processing chunk {idx+1}/{len(chunks)} (length: {len(chunk)})")
        
        # Adjust the system prompt for intermediary chunks vs. final summary
        if len(chunks) > 1:
            if idx < len(chunks) - 1:
                system_prompt = f"""You are a document summarization expert tasked with creating an intermediate summary of a portion of a document. 
                This is chunk {idx+1} of {len(chunks)}. Focus on extracting the main points and key information only.
                Be concise but thorough. The final summary will be created by combining these intermediary summaries.
                IMPORTANT: Do NOT truncate or omit ANY information from the document. Your summary must be comprehensive and include ALL key points.
                You MUST include at least 5-10 explicit key points in your response, formatted as a list."""
            else:
                system_prompt = f"""You are a document summarization expert tasked with combining all previous summaries into a final coherent summary.
                This is the final chunk {idx+1} of {len(chunks)}. Create a well-structured final summary that captures the key points from all chunks.
                The length should be {summary_options.get('length', 'medium')} and the style should be {summary_options.get('style', 'concise')}.
                IMPORTANT: Do NOT truncate or omit ANY information from the document. Your summary must be comprehensive and include ALL key points.
                You MUST include at least 10-15 explicit key points in your response, formatted as a list."""
        else:
            # For documents that fit in a single chunk
            system_prompt = f"""You are a document summarization expert tasked with creating a comprehensive summary of a document.
            The summary should be {summary_options.get('length', 'medium')} in length and {summary_options.get('style', 'concise')} in style.
            IMPORTANT: Do NOT truncate or omit ANY information from the document. Your summary must be comprehensive and include ALL key points and important details.
            You MUST include at least 10-15 explicit key points in your response, formatted as a list.
            The document has {total_pages} pages total."""
        
        # Prepare specific instructions based on document type and options
        specific_instructions = ""
        if summary_options.get('document_type') == 'pdf':
            specific_instructions = """For this PDF document, identify the main sections, key arguments, and supporting evidence."""
        elif summary_options.get('document_type') == 'pptx':
            specific_instructions = """For this presentation, focus on the main message of each slide and the overall narrative flow."""
        elif summary_options.get('document_type') == 'docx':
            specific_instructions = """For this text document, identify the thesis, main arguments, and supporting details."""
        elif summary_options.get('document_type') == 'xlsx':
            specific_instructions = """For this spreadsheet data, identify patterns, key metrics, and insights from the numeric data."""
        
        # Add structure instructions if requested
        structure_instructions = ""
        if summary_options.get('include_structure', False):
            structure_instructions = """In addition to the summary, provide an outline of the document's structure, identifying main sections and subsections."""
        
        # Create full prompt with all instructions
        full_system_prompt = f"""{system_prompt}
        {specific_instructions}
        {structure_instructions}
        
        Respond in JSON format with the following structure:
        {{
            "summary": "Comprehensive summary of the text",
            "key_points": ["Key point 1", "Key point 2", ...],
            "document_structure": {{"section_name": "description", ...}}
        }}
        """
        
        # Add retry logic with exponential backoff
        max_retries = 3
        retry_count = 0
        success = False
        response = None
        
        while retry_count < max_retries and not success:
            try:
                # Adjust chunk size if previous attempts failed
                if retry_count > 0:
                    logger.warning(f"Retry {retry_count}/{max_retries} - Reducing content size for reliability")
                    # Progressively reduce the chunk size on retries
                    reduced_size = int(len(chunk) * (0.7 ** retry_count))
                    reduced_chunk = chunk[:reduced_size]
                    user_content = f"Here's the document content to summarize:\n\n{reduced_chunk}"
                    logger.info(f"Reduced chunk size from {len(chunk)} to {reduced_size} characters")
                else:
                    user_content = f"Here's the document content to summarize:\n\n{chunk}"
                
                # Use gpt4o_mini_service or gpt4o_service based on content size
                # For larger documents or final summary, use the more powerful model
                if len(chunk) > 6000 or (len(chunks) > 1 and idx == len(chunks) - 1):
                    logger.info(f"Using gpt4o_service for complex summarization (attempt {retry_count+1}/{max_retries})")
                    response = gpt4o_service(
                        system_prompt=full_system_prompt.strip(),
                        user_input=user_content,
                        temperature=summary_options.get('temperature', 0.3),
                        json_output=True
                    )
                    if response.get("success", False) and "model" in response:
                        if response["model"] not in models_used:
                            models_used.append(response["model"])
                else:
                    logger.info(f"Using gpt4o_mini_service for summarization (attempt {retry_count+1}/{max_retries})")
                    response = gpt4o_mini_service(
                        system_prompt=full_system_prompt.strip(),
                        user_input=user_content,
                        temperature=summary_options.get('temperature', 0.3),
                        json_output=True
                    )
                    if response.get("success", False) and "model" in response:
                        if response["model"] not in models_used:
                            models_used.append(response["model"])
                
                if response.get("success", False):
                    success = True
                    logger.info(f"Successfully processed chunk {idx+1} on attempt {retry_count+1}")
                else:
                    logger.warning(f"LLM API call failed: {response.get('error', 'Unknown error')}")
                    retry_count += 1
                    
                    if retry_count < max_retries:
                        # Exponential backoff
                        sleep_time = 2 ** retry_count
                        logger.info(f"Waiting {sleep_time} seconds before retry...")
                        time.sleep(sleep_time)
                    else:
                        logger.error(f"Failed to process chunk after {max_retries} attempts")
                        # Create a minimal response to continue processing
                        response = {
                            "success": True,
                            "result": json.dumps({
                                "summary": f"[Processing error: Unable to summarize chunk {idx+1} after multiple attempts.]",
                                "key_points": ["Error occurred during processing"],
                                "document_structure": {}
                            }),
                            "input_tokens": 0,
                            "completion_tokens": 0,
                            "output_tokens": 0,
                            "model": "unknown"
                        }
                        if "unknown" not in models_used:
                            models_used.append("unknown")
            
            except Exception as e:
                logger.error(f"Exception during LLM API call: {str(e)}")
                retry_count += 1
                
                if retry_count < max_retries:
                    # Exponential backoff
                    sleep_time = 2 ** retry_count
                    logger.info(f"Waiting {sleep_time} seconds before retry...")
                    time.sleep(sleep_time)
                else:
                    logger.error(f"Failed to process chunk after {max_retries} attempts due to exception: {str(e)}")
                    # Create a minimal response to continue processing
                    response = {
                        "success": True,
                        "result": json.dumps({
                            "summary": f"[Error: Unable to process this section of the document. Technical error: {str(e)}]",
                            "key_points": ["Error occurred during processing"],
                            "document_structure": {}
                        }),
                        "input_tokens": 0,
                        "completion_tokens": 0,
                        "output_tokens": 0,
                        "model": "unknown"
                    }
                    if "unknown" not in models_used:
                        models_used.append("unknown")
        
        if not response or not response.get("success", False):
            logger.error(f"Error from LLM service: {response.get('error', 'Unknown error') if response else 'No response'}")
            # Create a fallback summary if API request fails
            all_summaries.append({
                "summary": f"Error: Unable to process document chunk {idx+1}: {response.get('error', 'Unknown error') if response else 'No response'}",
                "key_points": ["Error occurred during processing"],
                "document_structure": {},
                "model": "unknown"
            })
            if "unknown" not in models_used:
                models_used.append("unknown")
            continue
        
        # Extract token usage information
        token_usage["prompt_tokens"] += response.get("input_tokens", 0)
        token_usage["completion_tokens"] += response.get("completion_tokens", 0)
        token_usage["total_tokens"] += response.get("output_tokens", 0)
        logger.info(f"Token usage for chunk {idx+1}: {response.get('output_tokens', 0)} tokens")
        
        # Parse the LLM response into structured data
        try:
            result_text = response["result"]
            
            # Handle different response formats that might be returned by the LLM
            try:
                # Try to parse as JSON
                summary_result = json.loads(result_text)
                logger.info(f"Successfully parsed response as JSON")
            except json.JSONDecodeError:
                logger.warning(f"Failed to parse response as JSON, using raw text")
                summary_result = {
                    "summary": result_text,
                    "key_points": extract_key_points_from_content(result_text),
                    "document_structure": {}
                }
            
            # Ensure the summary has all required fields
            if not isinstance(summary_result, dict):
                logger.warning(f"Summary result is not a dictionary: {type(summary_result)}")
                summary_result = {
                    "summary": str(summary_result),
                    "key_points": [],
                    "document_structure": {}
                }
            
            if "summary" not in summary_result or not summary_result["summary"]:
                logger.warning("Missing or empty summary field in result")
                # Extract a summary from the content if possible
                if isinstance(result_text, str) and len(result_text) > 0:
                    summary_result["summary"] = result_text
                else:
                    summary_result["summary"] = f"Summary could not be generated for chunk {idx+1}."
            
            if "key_points" not in summary_result or not summary_result["key_points"]:
                logger.warning("No key points found in LLM response, extracting from content")
                # Try to extract key points from the summary
                summary_text = summary_result.get("summary", "")
                summary_result["key_points"] = extract_key_points_from_content(summary_text)
            
            if "document_structure" not in summary_result:
                summary_result["document_structure"] = {}
                
            # Add model information from response
            summary_result["model"] = response.get("model", "unknown")
            
            all_summaries.append(summary_result)
            logger.info(f"Successfully processed chunk {idx+1}, summary length: {len(summary_result['summary'])}")
            
        except Exception as e:
            logger.error(f"Error processing LLM response: {str(e)}")
            all_summaries.append({
                "summary": f"Error processing document chunk {idx+1}: {str(e)}",
                "key_points": ["Error occurred during processing"],
                "document_structure": {},
                "model": response.get("model", "unknown")
            })
    
    # If no summaries were generated at all, create a fallback
    if not all_summaries:
        logger.error("No summaries were generated for any chunks")
        return {
            "summary": "Failed to generate summary. The document could not be processed successfully.",
            "key_points": ["Document processing failed"],
            "document_structure": {},
            "pages_processed": total_pages,
            "tokens": token_usage,
            "model": "unknown"
        }
    
    # For multiple chunks, we need to combine them
    if len(all_summaries) > 1:
        # Extract all summaries and key points
        combined_text = "\n\n".join([s.get("summary", "") for s in all_summaries])
        all_key_points = []
        for summary in all_summaries:
            all_key_points.extend(summary.get("key_points", []))
        
        # Combine document structures
        combined_structure = {}
        for summary in all_summaries:
            combined_structure.update(summary.get("document_structure", {}))
        
        # Create final combined summary
        if len(combined_text) > 12000:
            # If combined text is still too long, create a final summary using LLM
            
            # Create full prompt with all instructions
            final_system_prompt = f"""You are a document summarization expert tasked with creating a final cohesive summary from multiple partial summaries.
            Create a well-structured final summary that integrates all the information.
            The length should be {summary_options.get('length', 'medium')} and the style should be {summary_options.get('style', 'concise')}.
            IMPORTANT: Do NOT truncate or omit ANY information from the document. Your summary must be comprehensive and include ALL key points.
            CRITICAL: You MUST extract and list at least 10-15 key points from the document as bullet points.
            
            Respond in JSON format with the following structure:
            {{
                "summary": "Final comprehensive summary",
                "key_points": ["Key point 1", "Key point 2", ...],
                "document_structure": {{"section_name": "description", ...}}
            }}
            """
            
            combined_content = f"Here are all the partial summaries to combine:\n\n{combined_text}\n\nAnd all key points:\n\n{json.dumps(all_key_points)}"
            
            # Add retry logic for final summary generation
            max_retries = 2
            for retry_attempt in range(max_retries):
                try:
                    # Use gpt4o_service for the final summary as it's more complex
                    response = gpt4o_service(
                        system_prompt=final_system_prompt,
                        user_input=combined_content,
                        temperature=summary_options.get('temperature', 0.3),
                        json_output=True
                    )
                    
                    if response["success"]:
                        if "model" in response and response["model"] not in models_used:
                            models_used.append(response["model"])
                        break
                    
                    logger.warning(f"Error generating final summary (attempt {retry_attempt+1}): {response.get('error', 'Unknown error')}")
                    
                    if retry_attempt < max_retries - 1:
                        # Reduce content size for retry
                        max_length = 10000 - (retry_attempt * 2000)  # Reduce by 2000 chars each retry
                        combined_content = f"Here are all the partial summaries to combine:\n\n{combined_text[:max_length]}\n\nAnd key points:\n\n{json.dumps(all_key_points[:min(15, len(all_key_points))])}"
                        time.sleep(2 ** retry_attempt)  # Exponential backoff
                    
                except Exception as e:
                    logger.error(f"Exception during final summary generation (attempt {retry_attempt+1}): {str(e)}")
                    if retry_attempt < max_retries - 1:
                        time.sleep(2 ** retry_attempt)  # Exponential backoff
            
            if not response["success"]:
                logger.error(f"Error generating final summary: {response.get('error', 'Unknown error')}")
                # Use the last summary as fallback but preserve ALL key points
                final_summary = all_summaries[-1]
                final_summary["key_points"] = all_key_points
                final_summary["document_structure"] = combined_structure
            else:
                # Add token usage
                token_usage["prompt_tokens"] += response.get("input_tokens", 0)
                token_usage["completion_tokens"] += response.get("completion_tokens", 0)
                token_usage["total_tokens"] += response.get("output_tokens", 0)
                
                # Parse response
                try:
                    result_text = response["result"]
                    
                    try:
                        # Try to parse as JSON
                        final_summary = json.loads(result_text)
                    except json.JSONDecodeError:
                        logger.error(f"Failed to parse JSON in final summary")
                        final_summary = {
                            "summary": result_text,
                            "key_points": all_key_points,
                            "document_structure": combined_structure
                        }
                    
                    # Add model information from response
                    final_summary["model"] = response.get("model", "unknown")
                    
                    # If the final summary doesn't include enough key points, check and add them
                    existing_key_points = final_summary.get("key_points", [])
                    if not existing_key_points:
                        logger.warning(f"Final summary contains no key points. Extracting from summary or using collected points.")
                        summary_text = final_summary.get("summary", "")
                        extracted_points = extract_key_points_from_content(summary_text)
                        
                        if extracted_points:
                            final_summary["key_points"] = extracted_points
                        else:
                            final_summary["key_points"] = all_key_points
                    elif len(existing_key_points) < min(5, len(all_key_points)):
                        logger.warning(f"Final summary contains too few key points ({len(existing_key_points)}). Adding all collected points.")
                        final_summary["key_points"] = all_key_points
                        
                except Exception as e:
                    logger.error(f"Error processing final summary: {str(e)}")
                    final_summary = {
                        "summary": combined_text[:2000] + "... [Summary truncated due to processing error]",
                        "key_points": all_key_points[:20],  # Include at most 20 key points
                        "document_structure": combined_structure,
                        "model": response.get("model", "unknown")
                    }
        else:
            # If combined text is manageable, use the last summary as the final one
            # but enrich it with all key points and structure
            final_summary = all_summaries[-1]
            final_summary["key_points"] = all_key_points
            final_summary["document_structure"] = combined_structure
    else:
        # For single chunk documents, just use the one summary
        final_summary = all_summaries[0]
    
    # Add page count and token usage to the final summary
    final_summary["pages_processed"] = total_pages
    final_summary["tokens"] = token_usage
    
    # Final validation to ensure we're returning usable data
    if not final_summary.get("summary"):
        logger.warning("Empty summary in final result, adding fallback text")
        final_summary["summary"] = "The document was processed, but no meaningful summary could be generated. Please check the document content or try again."
    
    # Ensure key points are present and valid
    if not final_summary.get("key_points") or len(final_summary.get("key_points", [])) == 0:
        logger.warning("No key points in final summary, generating from summary text")
        final_summary["key_points"] = extract_key_points_from_content(final_summary.get("summary", ""))
    
    # Make sure the summary is a string, not a dictionary
    if isinstance(final_summary["summary"], dict):
        logger.warning("Summary is a dictionary in final result, extracting text")
        summary_text, key_points, token_info = extract_summary_from_llm_response(final_summary["summary"])
        final_summary["summary"] = summary_text
        if key_points and not final_summary.get("key_points"):
            final_summary["key_points"] = key_points
        if token_info and (token_info["prompt_tokens"] > 0 or token_info["completion_tokens"] > 0 or token_info["total_tokens"] > 0):
            final_summary["tokens"] = token_info
            
    # Set model info in the final summary
    if "model" not in final_summary or not final_summary["model"]:
        # If multiple models were used, list them joined by a comma
        if models_used:
            final_summary["model"] = ", ".join(models_used)
        else:
            final_summary["model"] = "unknown"
    
    logger.info(f"Final summary generated with token usage: {token_usage}")
    return final_summary

def upload_summary_to_blob(summary_content, file_name, user_id, token):
    """
    Upload summary using FileService
    Args:
        summary_content: The summary content to upload
        file_name: Original file name to use as a base
        user_id: User ID for the request
        token: Authentication token
    Returns:
        File ID from the FileService
    """
    # Create a text file with the summary
    summary_file_name = f"summary_{os.path.splitext(file_name)[0]}.txt"
    
    # Convert summary to formatted text
    if isinstance(summary_content, dict):
        text_content = f"# Document Summary\n\n"
        text_content += f"## Executive Summary\n\n{summary_content.get('summary', '')}\n\n"
        
        text_content += "## Key Points\n\n"
        for idx, point in enumerate(summary_content.get('key_points', [])):
            text_content += f"{idx+1}. {point}\n"
        text_content += "\n"
        
        if summary_content.get('document_structure'):
            text_content += "## Document Structure\n\n"
            for section, description in summary_content.get('document_structure', {}).items():
                text_content += f"### {section}\n{description}\n\n"
        
        text_content += f"## Metadata\n\n"
        text_content += f"- Pages/Slides Processed: {summary_content.get('pages_processed', 0)}\n"
        text_content += f"- Tokens Used: {summary_content.get('tokens', {}).get('total_tokens', 0)}\n"
        text_content += f"  - Prompt Tokens: {summary_content.get('tokens', {}).get('prompt_tokens', 0)}\n"
        text_content += f"  - Completion Tokens: {summary_content.get('tokens', {}).get('completion_tokens', 0)}\n"
        text_content += f"- Model: {summary_content.get('model', 'unknown')}\n"
    else:
        # If it's not a dict, just use the content directly
        text_content = summary_content
    
    logger.info(f"Preparing to upload summary file: {summary_file_name}, size: {len(text_content)} characters")
    
    # For very large summaries, we may need to split them into multiple files
    max_upload_size = 10 * 1024 * 1024  # 10MB limit for upload
    
    if len(text_content) > max_upload_size:
        logger.warning(f"Summary text exceeds maximum upload size ({len(text_content)} > {max_upload_size})")
        # Split into multiple files if needed
        parts = []
        current_part = 1
        for i in range(0, len(text_content), max_upload_size):
            part_name = f"{summary_file_name}.part{current_part}"
            part_content = text_content[i:i+max_upload_size]
            parts.append((part_name, part_content))
            current_part += 1
            
        logger.info(f"Split summary into {len(parts)} parts for upload")
        
        # Upload each part
        uploaded_ids = []
        for part_name, part_content in parts:
            # Create a custom file-like object that mimics Flask's file object
            class MockFileObj:
                def __init__(self, data, filename, content_type):
                    self._io = io.BytesIO(data)
                    self.filename = filename
                    self.content_type = content_type
                
                def read(self):
                    return self._io.getvalue()
                
                def save(self, dst):
                    with open(dst, 'wb') as f:
                        f.write(self._io.getvalue())
            
            file_obj = MockFileObj(part_content.encode('utf-8'), part_name, 'text/plain')
            
            # Use FileService to upload the file
            file_info, error = FileService.upload_file(file_obj, user_id, FILE_UPLOAD_CONTAINER)
            
            if error:
                logger.error(f"Failed to upload summary file part {part_name}: {error}")
                raise Exception(f"Failed to upload summary file part {part_name}: {error}")
            
            uploaded_ids.append(file_info["file_id"])
        
        # Create a manifest file that lists all the parts
        manifest = f"# Summary Split into Multiple Files\n\nThis summary was split into {len(parts)} files due to size limitations.\n\n"
        for i, file_id in enumerate(uploaded_ids):
            manifest += f"Part {i+1}: {file_id}\n"
        
        # Upload the manifest
        manifest_name = f"{summary_file_name}.manifest"
        manifest_file = MockFileObj(manifest.encode('utf-8'), manifest_name, 'text/plain')
        
        manifest_info, error = FileService.upload_file(manifest_file, user_id, FILE_UPLOAD_CONTAINER)
        
        if error:
            logger.error(f"Failed to upload manifest file: {error}")
            # Return the first part ID if manifest upload fails
            return uploaded_ids[0] if uploaded_ids else None
        
        return manifest_info["file_id"]
    
    else:
        # Normal upload for reasonably sized files
        # Create a custom file-like object that mimics Flask's file object
        class MockFileObj:
            def __init__(self, data, filename, content_type):
                self._io = io.BytesIO(data)
                self.filename = filename
                self.content_type = content_type
            
            def read(self):
                return self._io.getvalue()
            
            def save(self, dst):
                with open(dst, 'wb') as f:
                    f.write(self._io.getvalue())
        
        file_obj = MockFileObj(text_content.encode('utf-8'), summary_file_name, 'text/plain')
        
        # Use FileService to upload the file
        file_info, error = FileService.upload_file(file_obj, user_id, FILE_UPLOAD_CONTAINER)
        
        if error:
            logger.error(f"Failed to upload summary file: {error}")
            raise Exception(f"Failed to upload summary file: {error}")
        
        return file_info["file_id"]

def document_summarization_route():
    """
    Summarize documents using AI (PDF, DOCX, PPTX, XLSX)
    ---
    tags:
      - Document Intelligence
    parameters:
      - name: X-Token
        in: header
        type: string
        required: true
        description: Valid token for authenticatio
      - name: X-Correlation-ID
        in: header
        type: string
        required: false
        description: Unique identifier for tracking requests across multiple systems
      - name: body
        in: body
        required: true
        schema:
          type: object
          required:
            - file_id
          properties:
            file_id:
              type: string
              description: ID of the uploaded file to summarize
            length:
              type: string
              enum: [short, medium, long, very_long]
              default: medium
              description: Desired summary length
            style:
              type: string
              enum: [concise, detailed, creative, technical, narrative, bullet_points]
              default: concise
              description: Style of the summary
            include_structure:
              type: boolean
              default: true
              description: Whether to include document structure in the summary
            temperature:
              type: number
              format: float
              minimum: 0
              maximum: 1.0
              default: 0.3
              description: Creativity level (0.0 = deterministic, 1.0 = creative)
            include_file_upload:
              type: boolean
              default: true
              description: Whether to create and upload a text file with the summary
    consumes:
      - application/json
    produces:
      - application/json
    responses:
      200:
        description: Document successfully summarized
        schema:
          type: object
          properties:
            message:
              type: string
              example: Document successfully summarized
            original_file_id:
              type: string
              description: ID of the original document
            summary_file_id:
              type: string
              description: ID of the uploaded summary file (if requested)
            summary:
              type: string
              description: Text summary of the document
            key_points:
              type: array
              items:
                type: string
              description: List of key points from the document
            document_structure:
              type: object
              description: Outline of document structure (if requested)
            pages_processed:
              type: integer
              description: Number of pages/slides/sheets processed
            token_usage:
              type: object
              properties:
                prompt_tokens:
                  type: integer
                completion_tokens:
                  type: integer
                total_tokens:
                  type: integer
              description: Token usage statistics
            model:
              type: string
              description: The AI model used for generating the summary
      400:
        description: Bad request
        schema:
          type: object
          properties:
            error:
              type: string
              example: Bad Request
            message:
              type: string
              example: Missing file_id parameter
      401:
        description: Authentication error
        schema:
          type: object
          properties:
            error:
              type: string
              example: Authentication Error
            message:
              type: string
              example: Missing X-Token header or Invalid token
      404:
        description: Not found
        schema:
          type: object
          properties:
            error:
              type: string
              example: Not Found
            message:
              type: string
              example: File with specified ID not found
      415:
        description: Unsupported media type
        schema:
          type: object
          properties:
            error:
              type: string
              example: Unsupported Media Type
            message:
              type: string
              example: File type not supported for summarization
      500:
        description: Server error
        schema:
          type: object
          properties:
            error:
              type: string
              example: Server Error
            message:
              type: string
              example: Error processing document
    """
    # Get token from X-Token header
    token = request.headers.get('X-Token')
    if not token:
        return create_api_response({
            "error": "Authentication Error",
            "message": "Missing X-Token header"
        }, 401)
    
    # Validate token
    token_details = DatabaseService.get_token_details_by_value(token)
    if not token_details:
        return create_api_response({
            "error": "Authentication Error",
            "message": "Invalid token"
        }, 401)
        
    # Check if token is expired
    now = datetime.now(pytz.UTC)
    expiration_time = token_details["token_expiration_time"]
    
    # Ensure expiration_time is timezone-aware
    if expiration_time.tzinfo is None:
        johannesburg_tz = pytz.timezone('Africa/Johannesburg')
        expiration_time = johannesburg_tz.localize(expiration_time)
        
    if now > expiration_time:
        return create_api_response({
            "error": "Authentication Error",
            "message": "Token has expired"
        }, 401)
        
    g.user_id = token_details["user_id"]
    g.token_id = token_details["id"]
    
    # Get request data
    data = request.get_json()
    if not data:
        return create_api_response({
            "error": "Bad Request",
            "message": "Request body is required"
        }, 400)
    
    # Validate file_id
    file_id = data.get('file_id')
    if not file_id:
        return create_api_response({
            "error": "Bad Request",
            "message": "file_id is required"
        }, 400)
    
    # Extract other parameters with defaults
    summary_options = {
        'length': data.get('length', 'medium'),
        'style': data.get('style', 'concise'),
        'include_structure': data.get('include_structure', True),
        'temperature': data.get('temperature', 0.3),
        'include_file_upload': data.get('include_file_upload', True),
        'token': token,  # Pass token for FileService
        'chunk_size': 8000  # Default chunk size for text processing - reduced for production
    }
    
    # Validate length parameter
    valid_lengths = ['short', 'medium', 'long', 'very_long']
    if summary_options['length'] not in valid_lengths:
        return create_api_response({
            "error": "Bad Request",
            "message": f"Invalid length parameter. Must be one of: {', '.join(valid_lengths)}"
        }, 400)
    
    # Validate style parameter
    valid_styles = ['concise', 'detailed', 'creative', 'technical', 'narrative', 'bullet_points']
    if summary_options['style'] not in valid_styles:
        return create_api_response({
            "error": "Bad Request",
            "message": f"Invalid style parameter. Must be one of: {', '.join(valid_styles)}"
        }, 400)
    
    # Validate temperature parameter
    if not 0 <= summary_options['temperature'] <= 1:
        return create_api_response({
            "error": "Bad Request",
            "message": "Temperature must be between 0.0 and 1.0"
        }, 400)
    
    temp_file_path = None
    try:
        # Get file URL using FileService instead of making an HTTP request
        file_info, error = FileService.get_file_url(file_id, g.user_id)
        
        if error:
            logger.error(f"Error retrieving file URL: {error}")
            return create_api_response({
                "error": "Not Found",
                "message": f"File with ID {file_id} not found or you don't have access"
            }, 404)
        
        file_url = file_info.get("file_url")
        file_name = file_info.get("file_name")
        
        logger.info(f"Retrieved file URL: {file_url} for file: {file_name}")
        
        if not file_url:
            logger.error("Missing file_url in file info")
            return create_api_response({
                "error": "Server Error",
                "message": "Failed to retrieve file URL"
            }, 500)
        
        # Download the file from the URL with better error handling
        try:
            logger.info(f"Downloading file from URL: {file_url}")
            
            # Try first with standard request
            file_response = requests.get(file_url, timeout=60)  # Increased timeout for larger files
            
            # If it's an Azure blob URL, we might need to add authorization
            if file_response.status_code != 200 and "blob.core.windows.net" in file_url:
                logger.info("Attempting download with authorization header")
                file_response = requests.get(
                    file_url,
                    headers={"Authorization": f"Bearer {token}"},
                    timeout=60
                )
            
            if file_response.status_code != 200:
                logger.error(f"Failed to download file: Status {file_response.status_code}")
                return create_api_response({
                    "error": "Server Error",
                    "message": f"Failed to download file from {file_url} (Status {file_response.status_code})"
                }, 500)
            
            content_length = len(file_response.content)
            logger.info(f"File downloaded successfully: {content_length} bytes")
            
            if content_length == 0:
                logger.error("Downloaded file is empty")
                return create_api_response({
                    "error": "Bad Request",
                    "message": "The downloaded file is empty"
                }, 400)
            
            # Save the file to a temporary location
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_name)[1]) as temp_file:
                temp_file.write(file_response.content)
                temp_file_path = temp_file.name
                
            logger.info(f"Downloaded file saved to temporary location: {temp_file_path}")
            
            # Verify the file exists and has content
            if not os.path.exists(temp_file_path) or os.path.getsize(temp_file_path) == 0:
                logger.error(f"Temporary file is missing or empty: {temp_file_path}")
                return create_api_response({
                    "error": "Server Error",
                    "message": "Downloaded file is empty or not accessible"
                }, 500)
                
        except Exception as e:
            logger.error(f"Error downloading file: {str(e)}")
            return create_api_response({
                "error": "Server Error",
                "message": f"Error downloading file: {str(e)}"
            }, 500)
        
        # Determine file type and extract text
        try:
            file_extension = os.path.splitext(file_name)[1].lower()
            
            logger.info(f"Processing file with extension: {file_extension}")
            
            if file_extension == '.pdf':
                logger.info(f"Extracting text from PDF file: {temp_file_path}")
                text_content, total_pages = extract_text_from_pdf(temp_file_path)
                summary_options['document_type'] = 'pdf'
            elif file_extension == '.docx':
                logger.info(f"Extracting text from DOCX file: {temp_file_path}")
                text_content, total_pages = extract_text_from_docx(temp_file_path)
                summary_options['document_type'] = 'docx'
            elif file_extension == '.pptx':
                logger.info(f"Extracting text from PPTX file: {temp_file_path}")
                text_content, total_pages = extract_text_from_pptx(temp_file_path)
                summary_options['document_type'] = 'pptx'
            elif file_extension in ['.xlsx', '.xls']:
                logger.info(f"Extracting text from Excel file: {temp_file_path}")
                text_content, total_pages = extract_text_from_xlsx(temp_file_path)
                summary_options['document_type'] = 'xlsx'
            else:
                # Clean up the temporary file
                if temp_file_path and os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
                    temp_file_path = None
                    
                logger.error(f"Unsupported file extension: {file_extension}")
                return create_api_response({
                    "error": "Unsupported Media Type",
                    "message": f"File type {file_extension} is not supported for summarization"
                }, 415)
            
            # Clean up the temporary file
            if temp_file_path and os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
                temp_file_path = None
            
            # Validate extracted text
            if not text_content or not text_content.strip():
                logger.error("No text content extracted from the document")
                return create_api_response({
                    "error": "Bad Request",
                    "message": "The document contains no extractable text content"
                }, 400)
                
            # Log the total content length for debugging
            content_length = len(text_content)
            logger.info(f"Document processed successfully: {content_length} characters extracted from {total_pages} pages")
            
            # Log a preview of the content
            if content_length > 0:
                preview = text_content[:200].replace('\n', ' ').strip()
                logger.info(f"Content preview: {preview}...")
            
        except Exception as e:
            # Clean up the temporary file if it exists
            if temp_file_path and os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
                temp_file_path = None
            
            logger.error(f"Error extracting text from file: {str(e)}")
            return create_api_response({
                "error": "Server Error",
                "message": f"Error extracting text from file: {str(e)}"
            }, 500)
            
        # Process text and generate summary using our chunk_and_summarize function
        # that now uses llmServices directly instead of making HTTP requests
        summary_result = chunk_and_summarize(text_content, total_pages, summary_options)
        
        # Format the summary based on document type
        if summary_options['document_type'] == 'pdf':
            formatted_summary = format_pdf_summary(summary_result)
        elif summary_options['document_type'] == 'docx':
            formatted_summary = format_docx_summary(summary_result)
        elif summary_options['document_type'] == 'pptx':
            formatted_summary = format_pptx_summary(summary_result)
        elif summary_options['document_type'] == 'xlsx':
            formatted_summary = format_xlsx_summary(summary_result)
        else:
            formatted_summary = summary_result
        
        # Upload summary file if requested
        summary_file_id = None
        if summary_options['include_file_upload']:
            try:
                summary_file_id = upload_summary_to_blob(formatted_summary, file_name, g.user_id, token)
            except Exception as e:
                logger.error(f"Error uploading summary file: {str(e)}")
                # Continue with the process even if file upload fails
        
        # Extract actual summary text, handling potential JSON wrapping
        summary_text = formatted_summary.get("summary", "")
        
        # Extract token usage directly from original response if available
        token_usage = formatted_summary.get("tokens", {})
        
        # Get model information
        model = formatted_summary.get("model", "unknown")
        
        # Handle case where summary is a JSON object
        if isinstance(summary_text, dict):
            logger.warning("Summary is a dictionary, extracting message field")
            
            # Extract token info if it exists in the summary dictionary
            if "input_tokens" in summary_text:
                token_usage["prompt_tokens"] = summary_text.get("input_tokens", 0)
            if "prompt_tokens" in summary_text:
                token_usage["prompt_tokens"] = summary_text.get("prompt_tokens", 0)
            if "completion_tokens" in summary_text:
                token_usage["completion_tokens"] = summary_text.get("completion_tokens", 0)
            if "output_tokens" in summary_text:
                token_usage["total_tokens"] = summary_text.get("output_tokens", 0)
            if "total_tokens" in summary_text:
                token_usage["total_tokens"] = summary_text.get("total_tokens", 0)
                
            # If it has a message field, that's the actual summary
            if "message" in summary_text:
                summary_text = summary_text["message"]
            # Otherwise stringify the whole object
            else:
                summary_text = str(summary_text)
                
        elif not isinstance(summary_text, str):
            logger.warning(f"Summary is not a string (type: {type(summary_text)}), converting")
            summary_text = str(summary_text)
        
        # Try to parse summary text if it looks like a JSON string
        if isinstance(summary_text, str) and summary_text.strip().startswith('{') and summary_text.strip().endswith('}'):
            try:
                summary_json = json.loads(summary_text)
                if isinstance(summary_json, dict):
                    # Extract token info if it exists in the parsed JSON
                    if "input_tokens" in summary_json:
                        token_usage["prompt_tokens"] = summary_json.get("input_tokens", 0)
                    if "prompt_tokens" in summary_json:
                        token_usage["prompt_tokens"] = summary_json.get("prompt_tokens", 0) 
                    if "completion_tokens" in summary_json:
                        token_usage["completion_tokens"] = summary_json.get("completion_tokens", 0)
                    if "output_tokens" in summary_json:
                        token_usage["total_tokens"] = summary_json.get("output_tokens", 0)
                    if "total_tokens" in summary_json:
                        token_usage["total_tokens"] = summary_json.get("total_tokens", 0)
                    
                    # Extract the actual summary text
                    if "message" in summary_json:
                        logger.info("Found message field in JSON string summary, extracting")
                        summary_text = summary_json["message"]
            except json.JSONDecodeError:
                # Not valid JSON, just use as is
                logger.info("Summary looks like JSON but isn't valid JSON, using as is")
                pass
        
        # Ensure key_points is a list of strings and not empty
        key_points = formatted_summary.get("key_points", [])
        if not key_points:
            logger.warning("No key points in formatted summary, generating from summary text")
            # Extract key points from the summary text
            key_points = extract_key_points_from_content(summary_text)
        
        # Ensure we have a valid token_usage object
        if not token_usage:
            logger.warning("No token usage information found, creating default")
            token_usage = {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0}
        
        # Fill in any missing token fields
        if "prompt_tokens" not in token_usage:
            token_usage["prompt_tokens"] = 0
        if "completion_tokens" not in token_usage:
            token_usage["completion_tokens"] = 0
        if "total_tokens" not in token_usage:
            # Calculate total if we have other values
            if token_usage["prompt_tokens"] > 0 or token_usage["completion_tokens"] > 0:
                token_usage["total_tokens"] = token_usage["prompt_tokens"] + token_usage["completion_tokens"]
            else:
                token_usage["total_tokens"] = 0
        
        logger.info(f"Token usage: {token_usage}")
        
        # Create response with proper structure
        response_data = {
            "message": "Document successfully summarized",
            "original_file_id": file_id,
            "summary_file_id": summary_file_id if summary_file_id else None,
            "summary": summary_text,
            "key_points": key_points,
            "token_usage": token_usage,  # Ensure token usage is included
            "model": model  # Include the model information
        }
        
        # Add the appropriate page/slide/sheet count
        if summary_options['document_type'] == 'pptx':
            response_data["slides_processed"] = formatted_summary.get("slides_processed", 0)
        elif summary_options['document_type'] == 'xlsx':
            response_data["sheets_processed"] = formatted_summary.get("sheets_processed", 0)
        else:
            response_data["pages_processed"] = formatted_summary.get("pages_processed", 0)
        
        # Add document structure if included
        if summary_options['include_structure']:
            if summary_options['document_type'] == 'pptx':
                response_data["presentation_structure"] = formatted_summary.get("presentation_structure", {})
            elif summary_options['document_type'] == 'xlsx':
                response_data["data_insights"] = formatted_summary.get("data_insights", {})
            else:
                response_data["document_structure"] = formatted_summary.get("document_structure", {})
                
        # Final verification - make sure summary is not JSON or dict
        if isinstance(response_data["summary"], dict):
            logger.error("Summary is still a dictionary after all processing!")
            response_data["summary"] = str(response_data["summary"])
        
        # Extract any wrapped JSON
        if isinstance(response_data["summary"], str) and response_data["summary"].strip().startswith('{') and response_data["summary"].strip().endswith('}'):
            try:
                json_obj = json.loads(response_data["summary"])
                if isinstance(json_obj, dict) and "message" in json_obj:
                    response_data["summary"] = json_obj["message"]
            except:
                pass
            
        logger.info(f"Final response data: {json.dumps(response_data, default=str)}")
        
        return create_api_response(response_data, 200)
        
    except Exception as e:
        # Clean up the temporary file if it exists
        if temp_file_path and os.path.exists(temp_file_path):
            os.unlink(temp_file_path)
        
        logger.error(f"Error in document summarization: {str(e)}")
        return create_api_response({
            "error": "Server Error",
            "message": f"Error processing document: {str(e)}"
        }, 500)

def register_document_intelligence_routes(app):
    """Register document intelligence routes with the Flask app"""
    app.route('/docint/summarization', methods=['POST'])(api_logger(check_balance(document_summarization_route)))
