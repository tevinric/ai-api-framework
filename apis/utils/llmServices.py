from openai import AzureOpenAI
from azure.ai.inference import ChatCompletionsClient
from azure.core.credentials import AzureKeyCredential
from apis.utils.config import (
    get_openai_client, 
    DEEPSEEK_API_KEY, 
    DEEPSEEK_V3_API_KEY, 
    LLAMA_API_KEY,
    O3_MINI_API_KEY,
    get_document_intelligence_config
)
import os
import tempfile
import base64
from azure.ai.formrecognizer import DocumentAnalysisClient
from apis.utils.fileService import FileService
import requests
import re
from collections import defaultdict
import math
import logging
import io
import time

# Configure logging
logger = logging.getLogger(__name__)

# Initialize OpenAI client
openai_client = get_openai_client()

# Initialize Document Intelligence client (used as fallback)
document_config = get_document_intelligence_config()
document_client = DocumentAnalysisClient(
    endpoint=document_config['endpoint'],
    credential=AzureKeyCredential(document_config['api_key'])
)

# Document types supported by Document Intelligence
DOCUMENT_INTELLIGENCE_FORMATS = {'pdf', 'doc', 'docx', 'xlsx', 'xls', 'pptx', 'ppt'}

# Allowed file extensions and MIME types
ALLOWED_EXTENSIONS = {
    'pdf': 'application/pdf',
    'txt': 'text/plain',
    'doc': 'application/msword',
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'csv': 'text/csv',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'xls': 'application/vnd.ms-excel',
    'ppt': 'application/vnd.ms-powerpoint',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg'
}

# Try importing document processing libraries
# These will be used first, falling back to Document Intelligence if they fail
NATIVE_LIBS_AVAILABLE = {}

try:
    import fitz  # PyMuPDF
    NATIVE_LIBS_AVAILABLE['pdf'] = True
    logger.info("PyMuPDF (fitz) library available for PDF processing")
except ImportError:
    NATIVE_LIBS_AVAILABLE['pdf'] = False
    logger.warning("PyMuPDF not available - will use Document Intelligence for PDFs")

try:
    import pandas as pd
    NATIVE_LIBS_AVAILABLE['excel'] = True
    logger.info("Pandas library available for Excel processing")
except ImportError:
    NATIVE_LIBS_AVAILABLE['excel'] = False
    logger.warning("Pandas not available - will use Document Intelligence for Excel files")

try:
    import docx
    NATIVE_LIBS_AVAILABLE['docx'] = True
    logger.info("python-docx library available for Word processing")
except ImportError:
    NATIVE_LIBS_AVAILABLE['docx'] = False
    logger.warning("python-docx not available - will use Document Intelligence for Word files")

try:
    from pptx import Presentation
    NATIVE_LIBS_AVAILABLE['pptx'] = True
    logger.info("python-pptx library available for PowerPoint processing")
except ImportError:
    NATIVE_LIBS_AVAILABLE['pptx'] = False
    logger.warning("python-pptx not available - will use Document Intelligence for PowerPoint files")

# Native document processing functions
def process_pdf_with_pymupdf(file_path):
    """Process PDF using PyMuPDF"""
    try:
        pdf_document = fitz.open(file_path)
        total_pages = len(pdf_document)
        text_content = ""
        toc_content = ""
        
        # Extract table of contents if available
        toc = pdf_document.get_toc()
        if toc:
            toc_content = "TABLE OF CONTENTS:\n"
            for entry in toc:
                level, title, page = entry
                indent = "  " * (level - 1)
                toc_content += f"{indent}* {title} (Page {page})\n"
            toc_content += "\n"
            
        # Extract text from each page
        for page_num in range(total_pages):
            page = pdf_document[page_num]
            text_content += f"\n--- Page {page_num + 1} ---\n"
            text_content += page.get_text()
            
        pdf_document.close()
        
        # Combine TOC and content
        full_content = f"PDF Document with {total_pages} pages\n\n"
        if toc_content:
            full_content += toc_content
        full_content += text_content
        
        return full_content, total_pages, True
    except Exception as e:
        logger.error(f"Error processing PDF with PyMuPDF: {str(e)}")
        return None, 0, False

def process_excel_with_pandas(file_path):
    """Process Excel files using pandas"""
    try:
        # Read all sheets
        xlsx = pd.ExcelFile(file_path)
        sheet_names = xlsx.sheet_names
        
        text_content = f"Excel Document with {len(sheet_names)} sheets\n\n"
        text_content += "SHEETS:\n"
        for sheet_name in sheet_names:
            text_content += f"* {sheet_name}\n"
        text_content += "\n"
        
        # Process each sheet
        for sheet_name in sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Handle empty sheets
            if df.empty:
                text_content += f"Sheet: {sheet_name} (empty)\n\n"
                continue
                
            text_content += f"Sheet: {sheet_name}\n"
            
            # Get column information
            text_content += f"Columns: {', '.join(df.columns)}\n"
            text_content += f"Rows: {len(df)}\n\n"
            
            # Show data sample (first 10 rows)
            max_rows = min(10, len(df))
            text_content += "Sample data:\n"
            text_content += df.head(max_rows).to_string(index=False) + "\n\n"
            
            # Add summary statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if not numeric_cols.empty:
                text_content += "Summary statistics for numeric columns:\n"
                text_content += df[numeric_cols].describe().to_string() + "\n\n"
                
        return text_content, len(sheet_names), True  # Count sheets as "pages"
    except Exception as e:
        logger.error(f"Error processing Excel with pandas: {str(e)}")
        return None, 0, False

def process_word_with_docx(file_path):
    """Process Word documents using python-docx"""
    try:
        doc = docx.Document(file_path)
        
        # Count paragraphs and approximate pages (rough estimate: 3000 chars per page)
        total_chars = sum(len(p.text) for p in doc.paragraphs)
        approx_pages = max(1, total_chars // 3000)
        
        text_content = f"Word Document (approximately {approx_pages} pages)\n\n"
        
        # Extract headers to build a TOC
        headers = []
        for paragraph in doc.paragraphs:
            if paragraph.style.name.startswith('Heading'):
                level = int(paragraph.style.name.replace('Heading', ''))
                headers.append((level, paragraph.text))
        
        if headers:
            text_content += "TABLE OF CONTENTS:\n"
            for level, title in headers:
                indent = "  " * (level - 1)
                text_content += f"{indent}* {title}\n"
            text_content += "\n"
        
        # Extract main content
        text_content += "DOCUMENT CONTENT:\n\n"
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text_content += paragraph.text + "\n\n"
                
        # Extract tables if any
        if doc.tables:
            text_content += "\nTABLES:\n"
            for i, table in enumerate(doc.tables):
                text_content += f"\nTable {i+1}:\n"
                
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    text_content += " | ".join(row_text) + "\n"
                text_content += "\n"
                
        return text_content, approx_pages, True
    except Exception as e:
        logger.error(f"Error processing Word document with python-docx: {str(e)}")
        return None, 0, False

def process_powerpoint_with_pptx(file_path):
    """Process PowerPoint presentations using python-pptx"""
    try:
        prs = Presentation(file_path)
        text_content = f"PowerPoint Presentation with {len(prs.slides)} slides\n\n"
        
        text_content += "SLIDES:\n"
        
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            text_content += f"\n--- Slide {slide_num} ---\n"
            
            # Extract title if available
            if slide.shapes.title:
                text_content += f"Title: {slide.shapes.title.text}\n\n"
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Skip the title that we've already included
                    if shape != slide.shapes.title:
                        text_content += f"{shape.text}\n"
            
            text_content += "\n"
            
        return text_content, len(prs.slides), True
    except Exception as e:
        logger.error(f"Error processing PowerPoint with python-pptx: {str(e)}")
        return None, 0, False

def process_document_with_native_libs(file_path, filename):
    """
    Process document with native libraries based on file type
    Returns: (content, page_count, used_native_lib)
    """
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    # PDF processing
    if ext == 'pdf' and NATIVE_LIBS_AVAILABLE['pdf']:
        return process_pdf_with_pymupdf(file_path)
        
    # Excel processing
    elif ext in ['xlsx', 'xls'] and NATIVE_LIBS_AVAILABLE['excel']:
        return process_excel_with_pandas(file_path)
        
    # Word processing
    elif ext in ['docx', 'doc'] and NATIVE_LIBS_AVAILABLE['docx']:
        if ext == 'doc':
            logger.warning("DOC format may not be fully supported by python-docx, might fall back to Document Intelligence")
        return process_word_with_docx(file_path)
        
    # PowerPoint processing
    elif ext in ['pptx', 'ppt'] and NATIVE_LIBS_AVAILABLE['pptx']:
        if ext == 'ppt':
            logger.warning("PPT format may not be fully supported by python-pptx, might fall back to Document Intelligence")
        return process_powerpoint_with_pptx(file_path)
        
    # No native library available for this file type
    return None, 0, False

def chunk_long_text(text, max_chunk_size=8000, overlap=200):
    """
    Chunk long text into overlapping segments that preserve context
    Try to break at paragraph boundaries when possible
    """
    if len(text) <= max_chunk_size:
        return [text]
        
    chunks = []
    
    # First try to split by double newlines (paragraphs)
    paragraphs = re.split(r'\n\s*\n', text)
    
    current_chunk = ""
    for paragraph in paragraphs:
        # If adding this paragraph would exceed max size, 
        # save current chunk and start a new one (with overlap)
        if len(current_chunk) + len(paragraph) > max_chunk_size:
            chunks.append(current_chunk)
            
            # Start new chunk with overlap from previous chunk if possible
            if len(current_chunk) > overlap:
                current_chunk = current_chunk[-overlap:] + "\n\n"
            else:
                current_chunk = ""
                
        # Add paragraph to current chunk
        if current_chunk and not current_chunk.endswith("\n\n"):
            current_chunk += "\n\n"
        current_chunk += paragraph
    
    # Add the last chunk if it has content
    if current_chunk:
        chunks.append(current_chunk)
    
    # If any chunks are still too large, split them further by sentences
    final_chunks = []
    for chunk in chunks:
        if len(chunk) <= max_chunk_size:
            final_chunks.append(chunk)
        else:
            # Split by sentences and rebuild chunks
            sentences = re.split(r'(?<=[.!?])\s+', chunk)
            sub_chunk = ""
            
            for sentence in sentences:
                if len(sub_chunk) + len(sentence) > max_chunk_size:
                    final_chunks.append(sub_chunk)
                    sub_chunk = ""
                
                if sub_chunk and not sub_chunk.endswith(" "):
                    sub_chunk += " "
                sub_chunk += sentence
            
            if sub_chunk:
                final_chunks.append(sub_chunk)
    
    return final_chunks

def deepseek_v3_service(system_prompt, user_input, temperature=0.7, json_output=False, max_tokens=1000):
    """DeepSeek-V3 LLM service function for general task completion"""
    try:
        # Fixed deployment endpoint
        ENDPOINT = 'https://ai-coe-services-dev.services.ai.azure.com/models'
        MODEL_NAME = 'DeepSeek-V3'
        
        # Initialize Azure Inference client
        client = ChatCompletionsClient(
            endpoint=ENDPOINT,
            credential=AzureKeyCredential(DEEPSEEK_V3_API_KEY)
        )
        
        # Prepare messages for the model
        messages = []
        
        # Add system message if provided
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        
        # Add user message
        messages.append({"role": "user", "content": user_input})
        
        # Prepare payload
        payload = {
            "messages": messages,
            "max_tokens": max_tokens,
            "temperature": temperature,
            "model": MODEL_NAME
        }
        
        # Add response format if JSON output is requested
        if json_output:
            payload["response_format"] = {"type": "json_object"}
        
        # Make request to LLM
        response = client.complete(payload)
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        model_name = response.model if hasattr(response, 'model') else MODEL_NAME
        cached_tokens = 0  # Default to 0 as this model doesn't support cached tokens
        
        return {
            "success": True,
            "result": result,
            "model": model_name,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens
        }
        
    except Exception as e:
        logger.error(f"DeepSeek-V3 API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def gpt4o_service(system_prompt, user_input, temperature=0.5, json_output=False):
    """OpenAI GPT-4o LLM service function for text completion and content generation"""
    try:
        # Fixed deployment model
        DEPLOYMENT = 'gpt-4o'
        
        # Make request to LLM
        response = openai_client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Provided text: {user_input}"}
            ],
            temperature=temperature,
            response_format={"type": "json_object"} if json_output else {"type": "text"}
        )
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
        
        return {
            "success": True,
            "result": result,
            "model": DEPLOYMENT,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens
        }
        
    except Exception as e:
        logger.error(f"GPT-4o API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def gpt4o_mini_service(system_prompt, user_input, temperature=0.5, json_output=False):
    """OpenAI GPT-4o-mini LLM service function for everyday text completion tasks"""
    try:
        # Fixed deployment model
        DEPLOYMENT = 'gpt-4o-mini'
        
        # Make request to LLM
        response = openai_client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Provided text: {user_input}"}
            ],
            temperature=temperature,
            response_format={"type": "json_object"} if json_output else {"type": "text"}
        )
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
        
        return {
            "success": True,
            "result": result,
            "model": DEPLOYMENT,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens
        }
        
    except Exception as e:
        logger.error(f"GPT-4o-mini API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def o1_mini_service(system_prompt, user_input, temperature=0.5, json_output=False):
    """OpenAI O1-mini LLM service function for complex tasks requiring reasoning"""
    try:
        # Fixed deployment model
        DEPLOYMENT = 'o1-mini'
        
        # o1-mini doesn't support 'system' role, so include system prompt in user message
        combined_input = f"{system_prompt}\n\nProvided text: {user_input}"
        
        # Make request to LLM
        response = openai_client.chat.completions.create(
            model=DEPLOYMENT,
            messages=[
                {"role": "user", "content": combined_input}
            ],
            response_format={"type": "json_object"} if json_output else {"type": "text"}
        )
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
        
        return {
            "success": True,
            "result": result,
            "model": DEPLOYMENT,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens
        }
        
    except Exception as e:
        logger.error(f"O1-mini API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def o3_mini_service(system_prompt, user_input, max_completion_tokens=100000, reasoning_effort="medium", json_output=False):
    """O3-Mini LLM service function with variable reasoning effort"""
    try:
        # Fixed deployment information
        ENDPOINT = "https://ai-coe-services-dev.openai.azure.com/"
        DEPLOYMENT = "o3-mini"
        API_VERSION = "2024-12-01-preview"
        
        # Initialize Azure OpenAI client
        client = AzureOpenAI(
            azure_endpoint=ENDPOINT,
            api_key=O3_MINI_API_KEY,
            api_version=API_VERSION,
        )
        
        # Prepare messages for the model
        chat_prompt = [
            {
                "role": "developer",
                "content": [
                    {
                        "type": "text",
                        "text": system_prompt
                    }
                ]
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": user_input
                    }
                ]
            }
        ]
        
        # Prepare additional parameters
        additional_params = {
            "max_completion_tokens": max_completion_tokens,
            "reasoning_effort": reasoning_effort
        }
        
        # Add response format if JSON output is requested
        if json_output:
            additional_params["response_format"] = {"type": "json_object"}
        
        # Make request to LLM
        completion = client.chat.completions.create(
            model=DEPLOYMENT,
            messages=chat_prompt,
            **additional_params,
            stop=None,
            stream=False
        )
        
        # Extract response data
        result = completion.choices[0].message.content
        prompt_tokens = completion.usage.prompt_tokens
        completion_tokens = completion.usage.completion_tokens
        total_tokens = prompt_tokens + completion_tokens
        model_name = DEPLOYMENT
        cached_tokens = completion.usage.cached_tokens if hasattr(completion.usage, 'cached_tokens') else 0
        
        return {
            "success": True,
            "result": result,
            "model": model_name,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens
        }
        
    except Exception as e:
        logger.error(f"O3-Mini API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def llama_service(system_prompt, user_input, temperature=0.8, json_output=False, max_tokens=2048, top_p=0.1, presence_penalty=0, frequency_penalty=0):
    """Meta Llama LLM service function for text generation"""
    try:
        # Fixed deployment endpoint
        ENDPOINT = 'https://Meta-Llama-3-1-405B-aiapis.eastus.models.ai.azure.com'
        
        # Initialize Azure Inference client
        client = ChatCompletionsClient(
            endpoint=ENDPOINT,
            credential=AzureKeyCredential(LLAMA_API_KEY)
        )
        
        # Prepare messages for the model
        messages = []
        
        # Add system message if provided
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        
        # Add user message
        messages.append({"role": "user", "content": user_input})
        
        # Prepare payload
        payload = {
            "messages": messages,
            "max_tokens": max_tokens,
            "temperature": temperature,
            "top_p": top_p,
            "presence_penalty": presence_penalty,
            "frequency_penalty": frequency_penalty
        }
        
        # Add response format if JSON output is requested
        if json_output:
            payload["response_format"] = {"type": "json_object"}
        
        # Make request to LLM
        response = client.complete(payload)
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        model_name = "llama-3-1-405b"
        cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
        
        return {
            "success": True,
            "result": result,
            "model": model_name,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens
        }
        
    except Exception as e:
        logger.error(f"Llama API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

# Enhanced Helper functions for document processing
def is_image_file(filename):
    """Check if the file is an image"""
    if '.' not in filename:
        return False
    ext = filename.rsplit('.', 1)[1].lower()
    return ext in ['png', 'jpg', 'jpeg']

def is_document_intelligence_supported(filename):
    """Check if file is supported by Document Intelligence"""
    if '.' not in filename:
        return False
    ext = filename.rsplit('.', 1)[1].lower()
    return ext in DOCUMENT_INTELLIGENCE_FORMATS

def process_text_file(file_path):
    """Process a plain text file"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error processing text file: {str(e)}")
        return f"[Error processing text file: {str(e)}]"

# Base class for multimodal model services
class MultimodalService:
    """Base class for all multimodal LLM services"""
    
    def __init__(self, model_name, max_files=10, max_token_limit=8000):
        self.model_name = model_name
        self.max_files = max_files
        self.max_token_limit = max_token_limit
    
    def process_document(self, file_path, file_name, user_id, user_query):
        """Process a document and return content and stats"""
        # Default implementation - override in subclasses if needed
        content, page_count, used_doc_intelligence = process_document_generic(
            file_path, file_name, user_query
        )
        
        file_stats = {
            "documents_processed": 1,
            "pages_processed": page_count if used_doc_intelligence else 0
        }
        
        return content, file_stats
    
    def handle_long_content(self, content, user_input):
        """Handle content that exceeds token limits"""
        # Default implementation chunks and adds context
        if len(content) > self.max_token_limit:
            chunks = chunk_long_text(content, self.max_token_limit)
            context = f"[This document has been split into {len(chunks)} parts due to its length. "
            context += f"I'll process each part and provide a comprehensive response.]\n\n"
            
            # Combine user input with document chunks
            processed_content = context + user_input + "\n\nDOCUMENT PART 1 OF " + str(len(chunks)) + ":\n" + chunks[0]
            
            # Add note about remaining chunks
            if len(chunks) > 1:
                processed_content += f"\n\n[Document continues with {len(chunks)-1} more parts that will be analyzed]"
                
            return processed_content
        
        return content
        
# Concrete implementation for GPT-4o
class GPT4oMultimodalService(MultimodalService):
    def __init__(self):
        super().__init__("gpt-4o", max_files=20, max_token_limit=100000)
    
    def process(self, system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
        """Process a multimodal request with GPT-4o"""
        try:
            temp_files = []  # Keep track of temporary files for cleanup
            
            # Track file processing statistics
            file_stats = {
                "documents_processed": 0,
                "images_processed": 0,
                "text_files_processed": 0,
                "pages_processed": 0  # Only counted for Document Intelligence
            }
            
            # Track processing errors to report back to user
            processing_errors = []
            
            # Prepare message content
            message_content = []
            
            # First add the user's input text
            message_content.append({
                "type": "text", 
                "text": user_input
            })
            
            # Process file_ids if any
            if file_ids and isinstance(file_ids, list) and len(file_ids) > 0:
                # Check for too many files
                if len(file_ids) > self.max_files:
                    return {
                        "success": False,
                        "error": f"Too many files. {self.model_name} can process a maximum of {self.max_files} files per request."
                    }
                
                for file_id in file_ids:
                    # Get file details using FileService
                    file_info, error = FileService.get_file_url(file_id, user_id)
                    
                    if error:
                        logger.error(f"Error retrieving file {file_id}: {error}")
                        continue
                    
                    # Check if we have the necessary file information
                    if not file_info or 'file_url' not in file_info or 'content_type' not in file_info:
                        logger.error(f"Incomplete file information for file {file_id}")
                        continue
                    
                    file_url = file_info['file_url']
                    content_type = file_info['content_type']
                    file_name = file_info.get('file_name', f"file_{file_id}")
                    
                    try:
                        # Create a temporary file to store the downloaded content
                        fd, temp_path = tempfile.mkstemp(suffix=f'.{file_name.rsplit(".", 1)[1].lower()}' if '.' in file_name else '')
                        temp_files.append(temp_path)
                        
                        # Download the file from Azure Blob Storage
                        import requests
                        response = requests.get(file_url)
                        response.raise_for_status()
                        
                        with os.fdopen(fd, 'wb') as tmp:
                            tmp.write(response.content)
                        
                        # Check file type and process accordingly
                        if content_type.startswith('image/'):
                            # For images, encode as base64 and add to message
                            file_stats["images_processed"] += 1
                            
                            with open(temp_path, 'rb') as img_file:
                                base64_image = base64.b64encode(img_file.read()).decode('utf-8')
                            
                            # Add as image content
                            message_content.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{content_type};base64,{base64_image}"
                                }
                            })
                            
                        elif is_document_intelligence_supported(file_name) or file_name.endswith('.txt'):
                            # Process document with native libs or Document Intelligence
                            try:
                                if file_name.endswith('.txt'):
                                    file_stats["text_files_processed"] += 1
                                    with open(temp_path, 'r', encoding='utf-8', errors='ignore') as f:
                                        extracted_content = f.read()
                                    used_doc_intelligence = False
                                else:
                                    # Use the generic document processor
                                    extracted_content, page_count, used_doc_intelligence = process_document_generic(
                                        temp_path, file_name, user_input
                                    )
                                    file_stats["documents_processed"] += 1
                                    # Only track pages for Document Intelligence
                                    if used_doc_intelligence:
                                        file_stats["pages_processed"] += page_count
                                
                                # Handle very long content by chunking if needed
                                if len(extracted_content) > self.max_token_limit:
                                    logger.info(f"Document {file_name} is very long, handling with chunking strategy")
                                    extracted_content = self.handle_long_content(extracted_content, user_input)
                                
                                # Add to message content
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                                })
                            except Exception as doc_error:
                                error_msg = f"Error processing document {file_name}: {str(doc_error)}"
                                logger.error(error_msg)
                                processing_errors.append(error_msg)
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\nError processing document {file_name}: {str(doc_error)}\n"
                                })
                            
                        else:
                            # For unsupported file types, just mention they were uploaded
                            message_content.append({
                                "type": "text",
                                "text": f"\n\nA file named '{file_name}' was uploaded, but its content type ({content_type}) is not supported for extraction."
                            })
                            
                    except requests.RequestException as e:
                        logger.error(f"Error downloading file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nFailed to download file {file_name}: {str(e)}"
                        })
                    except Exception as e:
                        logger.error(f"Error processing file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nError processing file {file_name}: {str(e)}"
                        })
            
            # Create the chat completion request
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": message_content}
            ]
            
            # Make the API call
            response = openai_client.chat.completions.create(
                model=self.model_name,
                messages=messages,
                temperature=temperature,
                max_tokens=4000,  # Add reasonable limit
                response_format={"type": "json_object"} if json_output else {"type": "text"}
            )
            
            # Extract response data
            result = response.choices[0].message.content
            prompt_tokens = response.usage.prompt_tokens
            completion_tokens = response.usage.completion_tokens
            total_tokens = response.usage.total_tokens
            cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
            
            # Total files processed
            total_files_processed = file_stats["documents_processed"] + file_stats["images_processed"] + file_stats["text_files_processed"]
            
            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logger.error(f"Error removing temporary file {temp_file}: {str(e)}")
            
            # Add any processing errors to the response
            if processing_errors:
                file_stats["processing_errors"] = processing_errors
            
            return {
                "success": True,
                "result": result,
                "model": self.model_name,
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": total_tokens,
                "cached_tokens": cached_tokens,
                "files_processed": total_files_processed,
                "file_processing_details": file_stats
            }
            
        except Exception as e:
            # Clean up temporary files in case of error
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except:
                    pass
                    
            logger.error(f"{self.model_name} Multimodal API error: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }

# Concrete implementation for GPT-4o-mini
class GPT4oMiniMultimodalService(MultimodalService):
    def __init__(self):
        super().__init__("gpt-4o-mini", max_files=10, max_token_limit=50000)
    
    def process(self, system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
        """Process a multimodal request with GPT-4o-mini"""
        try:
            temp_files = []  # Keep track of temporary files for cleanup
            
            # Track file processing statistics
            file_stats = {
                "documents_processed": 0,
                "images_processed": 0,
                "text_files_processed": 0,
                "pages_processed": 0  # Only counted for Document Intelligence
            }
            
            # Track processing errors to report back to user
            processing_errors = []
            
            # Prepare message content
            message_content = []
            
            # First add the user's input text
            message_content.append({
                "type": "text", 
                "text": user_input
            })
            
            # Process file_ids if any
            if file_ids and isinstance(file_ids, list) and len(file_ids) > 0:
                # Check for too many files
                if len(file_ids) > self.max_files:
                    return {
                        "success": False,
                        "error": f"Too many files. {self.model_name} can process a maximum of {self.max_files} files per request."
                    }
                
                for file_id in file_ids:
                    # Get file details using FileService
                    file_info, error = FileService.get_file_url(file_id, user_id)
                    
                    if error:
                        logger.error(f"Error retrieving file {file_id}: {error}")
                        continue
                    
                    # Check if we have the necessary file information
                    if not file_info or 'file_url' not in file_info or 'content_type' not in file_info:
                        logger.error(f"Incomplete file information for file {file_id}")
                        continue
                    
                    file_url = file_info['file_url']
                    content_type = file_info['content_type']
                    file_name = file_info.get('file_name', f"file_{file_id}")
                    
                    try:
                        # Create a temporary file to store the downloaded content
                        fd, temp_path = tempfile.mkstemp(suffix=f'.{file_name.rsplit(".", 1)[1].lower()}' if '.' in file_name else '')
                        temp_files.append(temp_path)
                        
                        # Download the file from Azure Blob Storage
                        import requests
                        response = requests.get(file_url)
                        response.raise_for_status()
                        
                        with os.fdopen(fd, 'wb') as tmp:
                            tmp.write(response.content)
                        
                        # Check file type and process accordingly
                        if content_type.startswith('image/'):
                            # For images, encode as base64 and add to message
                            file_stats["images_processed"] += 1
                            
                            with open(temp_path, 'rb') as img_file:
                                base64_image = base64.b64encode(img_file.read()).decode('utf-8')
                            
                            # Add as image content
                            message_content.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{content_type};base64,{base64_image}"
                                }
                            })
                            
                        elif is_document_intelligence_supported(file_name) or file_name.endswith('.txt'):
                            # Process document with native libs or Document Intelligence
                            try:
                                if file_name.endswith('.txt'):
                                    file_stats["text_files_processed"] += 1
                                    with open(temp_path, 'r', encoding='utf-8', errors='ignore') as f:
                                        extracted_content = f.read()
                                    used_doc_intelligence = False
                                else:
                                    # Use the generic document processor
                                    extracted_content, page_count, used_doc_intelligence = process_document_generic(
                                        temp_path, file_name, user_input
                                    )
                                    file_stats["documents_processed"] += 1
                                    # Only track pages for Document Intelligence
                                    if used_doc_intelligence:
                                        file_stats["pages_processed"] += page_count
                                
                                # Handle very long content by chunking if needed
                                if len(extracted_content) > self.max_token_limit:
                                    logger.info(f"Document {file_name} is very long, handling with chunking strategy")
                                    extracted_content = self.handle_long_content(extracted_content, user_input)
                                
                                # Add to message content
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                                })
                            except Exception as doc_error:
                                error_msg = f"Error processing document {file_name}: {str(doc_error)}"
                                logger.error(error_msg)
                                processing_errors.append(error_msg)
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\nError processing document {file_name}: {str(doc_error)}\n"
                                })
                        
                        else:
                            # For unsupported file types, just mention they were uploaded
                            message_content.append({
                                "type": "text",
                                "text": f"\n\nA file named '{file_name}' was uploaded, but its content type ({content_type}) is not supported for extraction."
                            })
                            
                    except requests.RequestException as e:
                        logger.error(f"Error downloading file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nFailed to download file {file_name}: {str(e)}"
                        })
                    except Exception as e:
                        logger.error(f"Error processing file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nError processing file {file_name}: {str(e)}"
                        })
            
            # Create the chat completion request
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": message_content}
            ]
            
            # Make the API call
            response = openai_client.chat.completions.create(
                model=self.model_name,
                messages=messages,
                temperature=temperature,
                max_tokens=2048,  # Lower token limit for mini
                response_format={"type": "json_object"} if json_output else {"type": "text"}
            )
            
            # Extract response data
            result = response.choices[0].message.content
            prompt_tokens = response.usage.prompt_tokens
            completion_tokens = response.usage.completion_tokens
            total_tokens = response.usage.total_tokens
            cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
            
            # Total files processed
            total_files_processed = file_stats["documents_processed"] + file_stats["images_processed"] + file_stats["text_files_processed"]
            
            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logger.error(f"Error removing temporary file {temp_file}: {str(e)}")
            
            # Add any processing errors to the response
            if processing_errors:
                file_stats["processing_errors"] = processing_errors
            
            return {
                "success": True,
                "result": result,
                "model": self.model_name,
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": total_tokens,
                "cached_tokens": cached_tokens,
                "files_processed": total_files_processed,
                "file_processing_details": file_stats
            }
            
        except Exception as e:
            # Clean up temporary files in case of error
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except:
                    pass
                    
            logger.error(f"{self.model_name} Multimodal API error: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
            
# Factory for multimodal services
class MultimodalServiceFactory:
    """Factory for creating multimodal services"""
    
    _services = {}
    
    @classmethod
    def register_service(cls, model_name, service_class):
        """Register a new multimodal service"""
        cls._services[model_name] = service_class
        logger.info(f"Registered multimodal service for model: {model_name}")
        
    @classmethod
    def get_service(cls, model_name):
        """Get multimodal service by model name"""
        service_class = cls._services.get(model_name)
        if not service_class:
            raise ValueError(f"No multimodal service registered for model: {model_name}")
        return service_class()

# Register GPT-4o and GPT-4o-mini services
MultimodalServiceFactory.register_service("gpt-4o", GPT4oMultimodalService)
MultimodalServiceFactory.register_service("gpt-4o-mini", GPT4oMiniMultimodalService)

# Updated multimodal service functions that use the factory
def gpt4o_multimodal_service(system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
    """OpenAI GPT-4o LLM service function for multimodal content generation with file support"""
    service = MultimodalServiceFactory.get_service("gpt-4o")
    return service.process(system_prompt, user_input, temperature, json_output, file_ids, user_id)

def gpt4o_mini_multimodal_service(system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
    """OpenAI GPT-4o-mini LLM service function for multimodal content generation with file support"""
    service = MultimodalServiceFactory.get_service("gpt-4o-mini")
    return service.process(system_prompt, user_input, temperature, json_output, file_ids, user_id)

# Concrete implementation for GPT-4o
class GPT4oMultimodalService(MultimodalService):
    def __init__(self):
        super().__init__("gpt-4o", max_files=20, max_token_limit=100000)
    
    def process(self, system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
        """Process a multimodal request with GPT-4o"""
        try:
            temp_files = []  # Keep track of temporary files for cleanup
            
            # Track file processing statistics
            file_stats = {
                "documents_processed": 0,
                "images_processed": 0,
                "text_files_processed": 0,
                "pages_processed": 0  # Only counted for Document Intelligence
            }
            
            # Track processing errors to report back to user
            processing_errors = []
            
            # Prepare message content
            message_content = []
            
            # First add the user's input text
            message_content.append({
                "type": "text", 
                "text": user_input
            })
            
            # Process file_ids if any
            if file_ids and isinstance(file_ids, list) and len(file_ids) > 0:
                # Check for too many files
                if len(file_ids) > self.max_files:
                    return {
                        "success": False,
                        "error": f"Too many files. {self.model_name} can process a maximum of {self.max_files} files per request."
                    }
                
                for file_id in file_ids:
                    # Get file details using FileService
                    file_info, error = FileService.get_file_url(file_id, user_id)
                    
                    if error:
                        logger.error(f"Error retrieving file {file_id}: {error}")
                        continue
                    
                    # Check if we have the necessary file information
                    if not file_info or 'file_url' not in file_info or 'content_type' not in file_info:
                        logger.error(f"Incomplete file information for file {file_id}")
                        continue
                    
                    file_url = file_info['file_url']
                    content_type = file_info['content_type']
                    file_name = file_info.get('file_name', f"file_{file_id}")
                    
                    try:
                        # Create a temporary file to store the downloaded content
                        fd, temp_path = tempfile.mkstemp(suffix=f'.{file_name.rsplit(".", 1)[1].lower()}' if '.' in file_name else '')
                        temp_files.append(temp_path)
                        
                        # Download the file from Azure Blob Storage
                        import requests
                        response = requests.get(file_url)
                        response.raise_for_status()
                        
                        with os.fdopen(fd, 'wb') as tmp:
                            tmp.write(response.content)
                        
                        # Check file type and process accordingly
                        if content_type.startswith('image/'):
                            # For images, encode as base64 and add to message
                            file_stats["images_processed"] += 1
                            
                            with open(temp_path, 'rb') as img_file:
                                base64_image = base64.b64encode(img_file.read()).decode('utf-8')
                            
                            # Add as image content
                            message_content.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{content_type};base64,{base64_image}"
                                }
                            })
                            
                        elif is_document_intelligence_supported(file_name) or file_name.endswith('.txt'):
                            # Process document with native libs or Document Intelligence
                            try:
                                if file_name.endswith('.txt'):
                                    file_stats["text_files_processed"] += 1
                                    with open(temp_path, 'r', encoding='utf-8', errors='ignore') as f:
                                        extracted_content = f.read()
                                    used_doc_intelligence = False
                                else:
                                    # Use the generic document processor
                                    extracted_content, page_count, used_doc_intelligence = process_document_generic(
                                        temp_path, file_name, user_input
                                    )
                                    file_stats["documents_processed"] += 1
                                    # Only track pages for Document Intelligence
                                    if used_doc_intelligence:
                                        file_stats["pages_processed"] += page_count
                                
                                # Handle very long content by chunking if needed
                                if len(extracted_content) > self.max_token_limit:
                                    logger.info(f"Document {file_name} is very long, handling with chunking strategy")
                                    extracted_content = self.handle_long_content(extracted_content, user_input)
                                
                                # Add to message content
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                                })
                            except Exception as doc_error:
                                error_msg = f"Error processing document {file_name}: {str(doc_error)}"
                                logger.error(error_msg)
                                processing_errors.append(error_msg)
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\nError processing document {file_name}: {str(doc_error)}\n"
                                })
                            
                        else:
                            # For unsupported file types, just mention they were uploaded
                            message_content.append({
                                "type": "text",
                                "text": f"\n\nA file named '{file_name}' was uploaded, but its content type ({content_type}) is not supported for extraction."
                            })
                            
                    except requests.RequestException as e:
                        logger.error(f"Error downloading file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nFailed to download file {file_name}: {str(e)}"
                        })
                    except Exception as e:
                        logger.error(f"Error processing file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nError processing file {file_name}: {str(e)}"
                        })
            
            # Create the chat completion request
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": message_content}
            ]
            
            # Make the API call
            response = openai_client.chat.completions.create(
                model=self.model_name,
                messages=messages,
                temperature=temperature,
                max_tokens=4000,  # Add reasonable limit
                response_format={"type": "json_object"} if json_output else {"type": "text"}
            )
            
            # Extract response data
            result = response.choices[0].message.content
            prompt_tokens = response.usage.prompt_tokens
            completion_tokens = response.usage.completion_tokens
            total_tokens = response.usage.total_tokens
            cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
            
            # Total files processed
            total_files_processed = file_stats["documents_processed"] + file_stats["images_processed"] + file_stats["text_files_processed"]
            
            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logger.error(f"Error removing temporary file {temp_file}: {str(e)}")
            
            # Add any processing errors to the response
            if processing_errors:
                file_stats["processing_errors"] = processing_errors
            
            return {
                "success": True,
                "result": result,
                "model": self.model_name,
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": total_tokens,
                "cached_tokens": cached_tokens,
                "files_processed": total_files_processed,
                "file_processing_details": file_stats
            }
            
        except Exception as e:
            # Clean up temporary files in case of error
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except:
                    pass
                    
            logger.error(f"{self.model_name} Multimodal API error: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }

# Concrete implementation for GPT-4o-mini
class GPT4oMiniMultimodalService(MultimodalService):
    def __init__(self):
        super().__init__("gpt-4o-mini", max_files=10, max_token_limit=50000)
    
    def process(self, system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
        """Process a multimodal request with GPT-4o-mini"""
        try:
            temp_files = []  # Keep track of temporary files for cleanup
            
            # Track file processing statistics
            file_stats = {
                "documents_processed": 0,
                "images_processed": 0,
                "text_files_processed": 0,
                "pages_processed": 0  # Only counted for Document Intelligence
            }
            
            # Track processing errors to report back to user
            processing_errors = []
            
            # Prepare message content
            message_content = []
            
            # First add the user's input text
            message_content.append({
                "type": "text", 
                "text": user_input
            })
            
            # Process file_ids if any
            if file_ids and isinstance(file_ids, list) and len(file_ids) > 0:
                # Check for too many files
                if len(file_ids) > self.max_files:
                    return {
                        "success": False,
                        "error": f"Too many files. {self.model_name} can process a maximum of {self.max_files} files per request."
                    }
                
                for file_id in file_ids:
                    # Get file details using FileService
                    file_info, error = FileService.get_file_url(file_id, user_id)
                    
                    if error:
                        logger.error(f"Error retrieving file {file_id}: {error}")
                        continue
                    
                    # Check if we have the necessary file information
                    if not file_info or 'file_url' not in file_info or 'content_type' not in file_info:
                        logger.error(f"Incomplete file information for file {file_id}")
                        continue
                    
                    file_url = file_info['file_url']
                    content_type = file_info['content_type']
                    file_name = file_info.get('file_name', f"file_{file_id}")
                    
                    try:
                        # Create a temporary file to store the downloaded content
                        fd, temp_path = tempfile.mkstemp(suffix=f'.{file_name.rsplit(".", 1)[1].lower()}' if '.' in file_name else '')
                        temp_files.append(temp_path)
                        
                        # Download the file from Azure Blob Storage
                        import requests
                        response = requests.get(file_url)
                        response.raise_for_status()
                        
                        with os.fdopen(fd, 'wb') as tmp:
                            tmp.write(response.content)
                        
                        # Check file type and process accordingly
                        if content_type.startswith('image/'):
                            # For images, encode as base64 and add to message
                            file_stats["images_processed"] += 1
                            
                            with open(temp_path, 'rb') as img_file:
                                base64_image = base64.b64encode(img_file.read()).decode('utf-8')
                            
                            # Add as image content
                            message_content.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{content_type};base64,{base64_image}"
                                }
                            })
                            
                        elif is_document_intelligence_supported(file_name) or file_name.endswith('.txt'):
                            # Process document with native libs or Document Intelligence
                            try:
                                if file_name.endswith('.txt'):
                                    file_stats["text_files_processed"] += 1
                                    with open(temp_path, 'r', encoding='utf-8', errors='ignore') as f:
                                        extracted_content = f.read()
                                    used_doc_intelligence = False
                                else:
                                    # Use the generic document processor
                                    extracted_content, page_count, used_doc_intelligence = process_document_generic(
                                        temp_path, file_name, user_input
                                    )
                                    file_stats["documents_processed"] += 1
                                    # Only track pages for Document Intelligence
                                    if used_doc_intelligence:
                                        file_stats["pages_processed"] += page_count
                                
                                # Handle very long content by chunking if needed
                                if len(extracted_content) > self.max_token_limit:
                                    logger.info(f"Document {file_name} is very long, handling with chunking strategy")
                                    extracted_content = self.handle_long_content(extracted_content, user_input)
                                
                                # Add to message content
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                                })
                            except Exception as doc_error:
                                error_msg = f"Error processing document {file_name}: {str(doc_error)}"
                                logger.error(error_msg)
                                processing_errors.append(error_msg)
                                message_content.append({
                                    "type": "text",
                                    "text": f"\n\nError processing document {file_name}: {str(doc_error)}\n"
                                })
                        
                        else:
                            # For unsupported file types, just mention they were uploaded
                            message_content.append({
                                "type": "text",
                                "text": f"\n\nA file named '{file_name}' was uploaded, but its content type ({content_type}) is not supported for extraction."
                            })
                            
                    except requests.RequestException as e:
                        logger.error(f"Error downloading file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nFailed to download file {file_name}: {str(e)}"
                        })
                    except Exception as e:
                        logger.error(f"Error processing file {file_id}: {str(e)}")
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nError processing file {file_name}: {str(e)}"
                        })
            
            # Create the chat completion request
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": message_content}
            ]
            
            # Make the API call
            response = openai_client.chat.completions.create(
                model=self.model_name,
                messages=messages,
                temperature=temperature,
                max_tokens=2048,  # Lower token limit for mini
                response_format={"type": "json_object"} if json_output else {"type": "text"}
            )
            
            # Extract response data
            result = response.choices[0].message.content
            prompt_tokens = response.usage.prompt_tokens
            completion_tokens = response.usage.completion_tokens
            total_tokens = response.usage.total_tokens
            cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
            
            # Total files processed
            total_files_processed = file_stats["documents_processed"] + file_stats["images_processed"] + file_stats["text_files_processed"]
            
            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logger.error(f"Error removing temporary file {temp_file}: {str(e)}")
            
            # Add any processing errors to the response
            if processing_errors:
                file_stats["processing_errors"] = processing_errors
            
            return {
                "success": True,
                "result": result,
                "model": self.model_name,
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": total_tokens,
                "cached_tokens": cached_tokens,
                "files_processed": total_files_processed,
                "file_processing_details": file_stats
            }
            
        except Exception as e:
            # Clean up temporary files in case of error
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except:
                    pass
                    
            logger.error(f"{self.model_name} Multimodal API error: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }

def split_into_semantic_chunks(content, max_chunk_size=1000):
    """Split content into logical chunks based on structure"""
    # This is a simple implementation - could be enhanced with ML-based chunking
    chunks = []
    current_chunk = ""
    current_chunk_size = 0
    
    # Split by paragraph boundaries
    paragraphs = re.split(r'\n\s*\n', content)
    
    for para in paragraphs:
        para_size = len(para.split())
        
        # If paragraph itself is too large, split it further
        if para_size > max_chunk_size:
            sentences = re.split(r'(?<=[.!?])\s+', para)
            current_sentence_group = ""
            
            for sentence in sentences:
                if len(current_sentence_group.split()) + len(sentence.split()) < max_chunk_size:
                    current_sentence_group += " " + sentence if current_sentence_group else sentence
                else:
                    if current_sentence_group:
                        chunks.append(current_sentence_group.strip())
                    current_sentence_group = sentence
            
            if current_sentence_group:
                chunks.append(current_sentence_group.strip())
        
        # If this paragraph plus the current chunk is too big, start a new chunk
        elif current_chunk_size + para_size > max_chunk_size:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = para
            current_chunk_size = para_size
        else:
            # Add to current chunk
            current_chunk += "\n\n" + para if current_chunk else para
            current_chunk_size += para_size
    
    # Add the last chunk if there is one
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks

def calculate_relevance_score(text, query):
    """Calculate relevance score between text and query"""
    # Simple keyword-based scoring
    # In a production environment, this would use embeddings and vector similarity
    query_terms = set(query.lower().split())
    text_lower = text.lower()
    
    # Count term frequency
    term_count = sum(1 for term in query_terms if term in text_lower)
    
    # Calculate base score based on term frequency
    base_score = term_count / max(len(query_terms), 1)
    
    # Bonus for title-like sections or sections that contain complete phrases from query
    if any(re.search(r'^\s*' + re.escape(term) + r'\s*[:.-]', text, re.IGNORECASE) for term in query_terms):
        base_score *= 1.5
    
    # Check for exact phrase matches
    for i in range(len(query_terms) - 1):
        phrase = ' '.join(list(query_terms)[i:i+2])
        if phrase in text_lower:
            base_score *= 1.2
    
    # Look for exact matches of the full query
    if query.lower() in text_lower:
        base_score *= 2
    
    return base_score

def rank_by_relevance(chunks, query):
    """Rank content chunks by relevance to the query"""
    ranked_chunks = []
    
    for i, chunk in enumerate(chunks):
        score = calculate_relevance_score(chunk, query)
        ranked_chunks.append({
            "index": i,
            "chunk": chunk,
            "score": score
        })
    
    # Sort by score (descending)
    ranked_chunks.sort(key=lambda x: x["score"], reverse=True)
    
    return ranked_chunks

def summarize_content(content, max_length=200):
    """Create a brief summary of content"""
    # In a real implementation, this would use the model itself for summarization
    # Here we're using a simple heuristic approach
    
    # Extract first few sentences (max 3)
    sentences = re.split(r'(?<=[.!?])\s+', content)
    if len(sentences) <= 3:
        summary = ' '.join(sentences)
    else:
        summary = ' '.join(sentences[:3]) + '...'
    
    # Truncate if still too long
    if len(summary) > max_length:
        summary = summary[:max_length-3] + '...'
        
    return summary

def extract_document_outline(document_structure):
    """Extract a readable document outline from the structure"""
    outline = ""
    
    # Add metadata if available
    metadata = next((item for item in document_structure["structure"] if item["type"] == "metadata"), None)
    if metadata:
        meta_content = metadata["content"]
        outline += "Document Information:\n"
        if meta_content.get("title"):
            outline += f"Title: {meta_content['title']}\n"
        if meta_content.get("author"):
            outline += f"Author: {meta_content['author']}\n"
        outline += f"Pages: {meta_content.get('pages', 'Unknown')}\n\n"
    
    # Check for Excel-specific entries
    sheet_entries = [item for item in document_structure["structure"] if item["type"] == "sheet"]
    if sheet_entries:
        outline += "Excel Content:\n"
        for sheet in sheet_entries:
            outline += f"* {sheet.get('name', f'Sheet/Table {sheet.get('sheet_idx', 0) + 1}')}\n"
        outline += "\n"
        return outline  # For Excel files, return early with just the sheet list
        
    # Check for PowerPoint-specific entries
    slide_entries = [item for item in document_structure["structure"] if item["type"] == "slide"]
    if slide_entries:
        outline += "PowerPoint Content:\n"
        for slide_idx, slide in enumerate(slide_entries):
            # Try to find a title in the headings
            slide_heading = next((h for h in document_structure["headings"] if h.get("slide_num") == slide_idx + 1), None)
            if slide_heading:
                outline += f"* Slide {slide_idx + 1}: {slide_heading['text']}\n"
            else:
                outline += f"* Slide {slide_idx + 1}\n"
        outline += "\n"
        return outline  # For PowerPoint files, return early with just the slide list
    
    # Standard Word/PDF document processing
    # Add table of contents based on headings
    if document_structure["headings"]:
        outline += "Table of Contents:\n"
        for heading in document_structure["headings"]:
            indent = "  " * (heading["level"] - 1)
            outline += f"{indent}* {heading['text']} (Page {heading['page']})\n"
        outline += "\n"
    
    # Count tables
    tables = [item for item in document_structure["structure"] if item["type"] == "table"]
    if tables:
        outline += f"Document contains {len(tables)} tables.\n\n"
    
    return outline

def process_with_document_intelligence(file_path, filename):
    """Process a document file using Azure Document Intelligence (fallback method)"""
    try:
        # Determine the best model to use based on file type
        ext = filename.rsplit('.', 1)[1].lower()
        
        # Select appropriate prebuilt model
        if ext in ['xlsx', 'xls', 'csv']:
            model = "prebuilt-layout"  # Good for tables and structured data
        else:
            model = "prebuilt-document"  # General document understanding
        
        # Process the document
        with open(file_path, "rb") as document:
            poller = document_client.begin_analyze_document(model, document)
            result = poller.result()
        
        # Extract content
        extracted_text = ""
        
        # Track page count
        page_count = 0
        
        # Add document metadata if available
        if hasattr(result, 'metadata') and result.metadata:
            page_count = result.metadata.page_count  # Store page count
            extracted_text += f"Document Metadata:\n"
            extracted_text += f"Pages: {result.metadata.page_count}\n"
            if hasattr(result.metadata, 'author') and result.metadata.author:
                extracted_text += f"Author: {result.metadata.author}\n"
            if hasattr(result.metadata, 'title') and result.metadata.title:
                extracted_text += f"Title: {result.metadata.title}\n"
            extracted_text += "\n"
        else:
            # If metadata doesn't have page count, count the pages in the result
            page_count = len(result.pages) if hasattr(result, 'pages') else 0
        
        # Get page-by-page content
        for page_idx, page in enumerate(result.pages):
            extracted_text += f"\n--- Page {page_idx + 1} ---\n"
            
            # Extract text from paragraphs if available
            if hasattr(page, 'paragraphs') and page.paragraphs:
                for para in page.paragraphs:
                    extracted_text += f"{para.content}\n"
            # Otherwise extract from lines
            else:
                for line in page.lines:
                    extracted_text += f"{line.content}\n"
        
        # Extract tables if present
        if hasattr(result, 'tables') and result.tables:
            extracted_text += "\n--- Tables ---\n"
            for i, table in enumerate(result.tables):
                extracted_text += f"\nTable {i+1}:\n"
                
                # Build a text representation of the table
                prev_row_idx = -1
                for cell in table.cells:
                    if cell.row_index > prev_row_idx:
                        extracted_text += "\n"
                        prev_row_idx = cell.row_index
                    
                    extracted_text += f"{cell.content}\t"
                extracted_text += "\n"
        
        # Include key-value pairs if detected
        if hasattr(result, 'key_value_pairs') and result.key_value_pairs:
            extracted_text += "\n--- Key-Value Pairs ---\n"
            for pair in result.key_value_pairs:
                key = pair.key.content if pair.key else "N/A"
                value = pair.value.content if pair.value else "N/A"
                extracted_text += f"{key}: {value}\n"
                
        return extracted_text, page_count
    
    except Exception as e:
        logger.error(f"Error processing document with Document Intelligence: {str(e)}")
        return f"[Error processing document with Document Intelligence: {str(e)}]", 0

def process_with_intelligent_extraction(file_path, filename, user_query):
    """Process document with intelligent context extraction based on user query"""
    try:
        # Track page count
        page_count = 0
        
        # Determine the best model to use based on file type
        ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
        
        # Select appropriate prebuilt model
        if ext in ['xlsx', 'xls', 'csv']:
            model = "prebuilt-layout"  # Good for tables and structured data
        else:
            model = "prebuilt-document"  # General document understanding
            
        logger.info(f"Processing {filename} with model: {model}")
            
        # Extract full document first
        with open(file_path, "rb") as document:
            poller = document_client.begin_analyze_document(model, document)
            result = poller.result()
        
        # Get page count from metadata if available
        if hasattr(result, 'metadata') and result.metadata:
            page_count = result.metadata.page_count
        else:
            # If metadata doesn't have page count, count the pages in the result
            page_count = len(result.pages) if hasattr(result, 'pages') else 0
            
        # Extract document structure - pass the file extension for special handling
        doc_structure = extract_document_structure(result, ext)
        
        # Extract full raw content
        raw_content = ""
        
        # Special handling for Excel files
        if ext in ['xlsx', 'xls', 'csv']:
            raw_content += f"--- Excel File: {filename} ---\n\n"
            
            # Extract tables content
            if hasattr(result, 'tables') and result.tables:
                for i, table in enumerate(result.tables):
                    raw_content += f"\nTable/Sheet {i+1}:\n"
                    
                    # Build a text representation of the table
                    prev_row_idx = -1
                    for cell in table.cells:
                        if cell.row_index > prev_row_idx:
                            raw_content += "\n"
                            prev_row_idx = cell.row_index
                        
                        raw_content += f"{cell.content}\t"
                    raw_content += "\n"
            
            # If no page count set, use 1 as minimum
            if page_count == 0:
                page_count = 1
        
        # Special handling for PowerPoint files
        elif ext in ['pptx', 'ppt']:
            raw_content += f"--- PowerPoint Presentation: {filename} ---\n\n"
            
            # Get slide-by-slide content
            for page_idx, page in enumerate(result.pages):
                raw_content += f"\n--- Slide {page_idx + 1} ---\n"
                
                # Use first line as potential title
                if hasattr(page, 'lines') and page.lines and len(page.lines) > 0:
                    raw_content += f"Title: {page.lines[0].content}\n\n"
                
                # Extract text from paragraphs if available
                if hasattr(page, 'paragraphs') and page.paragraphs:
                    for para in page.paragraphs:
                        raw_content += f"{para.content}\n\n"
                # Otherwise extract from lines (skipping first line if used as title)
                elif hasattr(page, 'lines') and page.lines:
                    start_idx = 1 if len(page.lines) > 0 else 0  # Skip title line
                    for line_idx in range(start_idx, len(page.lines)):
                        raw_content += f"{page.lines[line_idx].content}\n"
            
            # If no page count set, use slide count
            if page_count == 0 and hasattr(result, 'pages'):
                page_count = len(result.pages)
        
        # Default handling for Word/PDF files
        else:
            # Get page-by-page content
            for page_idx, page in enumerate(result.pages):
                # Extract text from paragraphs if available
                if hasattr(page, 'paragraphs') and page.paragraphs:
                    for para in page.paragraphs:
                        raw_content += f"{para.content}\n\n"
                # Otherwise extract from lines
                else:
                    for line in page.lines:
                        raw_content += f"{line.content}\n"
        
        # Extract tables if present (for Word/PDF)
        tables_content = ""
        if ext not in ['xlsx', 'xls', 'csv'] and hasattr(result, 'tables') and result.tables:
            for i, table in enumerate(result.tables):
                tables_content += f"\nTable {i+1}:\n"
                
                # Build a text representation of the table
                prev_row_idx = -1
                for cell in table.cells:
                    if cell.row_index > prev_row_idx:
                        tables_content += "\n"
                        prev_row_idx = cell.row_index
                    
                    tables_content += f"{cell.content}\t"
                tables_content += "\n"
        
        # Add tables to the content for Word/PDF
        if tables_content and ext not in ['xlsx', 'xls', 'csv']:
            raw_content += "\n--- Tables ---\n" + tables_content
            
        # Split content into semantic chunks
        chunks = split_into_semantic_chunks(raw_content)
        
        # Rank chunks by relevance
        ranked_chunks = rank_by_relevance(chunks, user_query)
        
        # Build final content
        final_content = f"--- Document: {filename} ({ext.upper()} file) ---\n\n"
        
        # Add document structure
        doc_outline = extract_document_outline(doc_structure)
        final_content += f"DOCUMENT OUTLINE:\n{doc_outline}\n\n"
        
        # Calculate token budget to dynamically adjust content inclusion
        # Simple approximation: 4 characters ≈ 1 token
        outline_tokens = len(doc_outline) // 4  
        
        # Allocate tokens for most relevant content, summaries, etc.
        # For GPT-4o, we'll target staying under 100K total tokens
        MAX_CONTEXT_TOKENS = 90000  # Conservative limit
        
        # Reserve tokens for top chunks and summary
        remaining_tokens = MAX_CONTEXT_TOKENS - outline_tokens - 1000  # Reserve 1000 tokens for overhead
        
        # Add most relevant chunks
        final_content += "MOST RELEVANT CONTENT:\n"
        
        tokens_used = 0
        chunks_included = 0
        
        for chunk_info in ranked_chunks:
            chunk = chunk_info["chunk"]
            chunk_tokens = len(chunk) // 4
            
            if tokens_used + chunk_tokens <= remaining_tokens * 0.7:  # Use 70% for most relevant chunks
                final_content += f"\n[Content Chunk {chunks_included + 1} - Relevance Score: {chunk_info['score']:.2f}]\n{chunk}\n"
                tokens_used += chunk_tokens
                chunks_included += 1
            else:
                # If we can't include the full chunk, include a summary
                summary = summarize_content(chunk)
                summary_tokens = len(summary) // 4
                
                if tokens_used + summary_tokens <= remaining_tokens * 0.7:
                    final_content += f"\n[Content Chunk {chunks_included + 1} - Summary - Relevance Score: {chunk_info['score']:.2f}]\n{summary}\n"
                    tokens_used += summary_tokens
                    chunks_included += 1
                else:
                    # Stop adding content if we're approaching the token limit
                    break
        
        # Add summary of remaining content
        if chunks_included < len(ranked_chunks):
            final_content += f"\n\nSUMMARY OF ADDITIONAL CONTENT:\n"
            final_content += f"There are {len(ranked_chunks) - chunks_included} additional content chunks that were not included in full due to length constraints.\n"
            
            # Add brief summaries of some remaining chunks
            remaining_token_budget = remaining_tokens - tokens_used
            summaries_to_add = min(5, len(ranked_chunks) - chunks_included)  # Add up to 5 summaries
            
            for i in range(chunks_included, chunks_included + summaries_to_add):
                if i < len(ranked_chunks):
                    chunk = ranked_chunks[i]["chunk"]
                    summary = summarize_content(chunk, max_length=100)  # Shorter summary
                    summary_tokens = len(summary) // 4
                    
                    if remaining_token_budget >= summary_tokens:
                        final_content += f"\n- {summary} [Relevance Score: {ranked_chunks[i]['score']:.2f}]"
                        remaining_token_budget -= summary_tokens
                    else:
                        break
        
        return final_content, page_count
    
    except Exception as e:
        logger.error(f"Error in intelligent document processing: {str(e)}")
        # Fall back to basic document processing if advanced processing fails
        return process_with_document_intelligence(file_path, filename)

def gpt4o_document_intelligence_service(system_prompt, user_input, temperature=0.5, files=None):
    """GPT-4o with Document Intelligence service function for file processing"""
    try:
        # Fixed deployment model
        DEPLOYMENT = 'gpt-4o'
        temp_files = []  # Keep track of temporary files for cleanup
        
        # Track file processing statistics
        file_stats = {
            "documents_processed": 0,
            "images_processed": 0,
            "text_files_processed": 0,
            "pages_processed": 0  # Add this field to track page counts
        }
        
        # Prepare message content
        message_content = []
        
        # First add the user's input text
        message_content.append({
            "type": "text", 
            "text": user_input
        })
        
        # Process uploaded files if any
        if files:
            for file in files:
                # Create a temporary file to store the uploaded content
                fd, temp_path = tempfile.mkstemp(suffix=f'.{file.filename.rsplit(".", 1)[1].lower()}')
                temp_files.append(temp_path)
                
                # Save the uploaded file to the temporary path
                with os.fdopen(fd, 'wb') as tmp:
                    file.save(tmp)
                
                # Check file type and process accordingly
                if is_image_file(file.filename):
                    # For images, encode as base64 and add to message
                    file_stats["images_processed"] += 1
                    
                    with open(temp_path, 'rb') as img_file:
                        base64_image = base64.b64encode(img_file.read()).decode('utf-8')
                    
                    # Get MIME type
                    ext = file.filename.rsplit('.', 1)[1].lower()
                    mime_type = ALLOWED_EXTENSIONS[ext]
                    
                    # Add as image content
                    message_content.append({
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{mime_type};base64,{base64_image}"
                        }
                    })
                    
                elif is_document_intelligence_supported(file.filename):
                    # Process document with intelligent extraction
                    file_stats["documents_processed"] += 1
                    
                    # Process document
                    extracted_content, page_count = process_with_intelligent_extraction(temp_path, file.filename, user_input)
                    file_stats["pages_processed"] += page_count  # Track page count
                    
                    # Add to message content
                    message_content.append({
                        "type": "text",
                        "text": f"\n\n--- Content from {file.filename} ---\n{extracted_content}\n--- End of {file.filename} ---\n"
                    })
                    
                elif file.filename.endswith('.txt'):
                    # Process plain text file
                    file_stats["text_files_processed"] += 1
                    
                    # Extract text
                    extracted_content = process_text_file(temp_path)
                    
                    # Add to message content
                    message_content.append({
                        "type": "text",
                        "text": f"\n\n--- Content from {file.filename} ---\n{extracted_content}\n--- End of {file.filename} ---\n"
                    })
                    
                else:
                    # For unsupported file types, just mention they were uploaded
                    message_content.append({
                        "type": "text",
                        "text": f"\n\nA file named '{file.filename}' was uploaded, but its content type is not supported for extraction."
                    })
        
        # Create the chat completion request
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": message_content}
        ]
        
        # Make the API call
        response = openai_client.chat.completions.create(
            model=DEPLOYMENT,
            messages=messages,
            temperature=temperature,
            max_tokens=4000  # Add reasonable limit
        )
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
        
        # Total files processed
        total_files_processed = file_stats["documents_processed"] + file_stats["images_processed"] + file_stats["text_files_processed"]
        
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except Exception as e:
                logger.error(f"Error removing temporary file {temp_file}: {str(e)}")
        
        # Add any processing errors to the response
        if processing_errors:
            file_stats["processing_errors"] = processing_errors
        
        return {
            "success": True,
            "result": result,
            "model": DEPLOYMENT,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens,
            "files_processed": total_files_processed,
            "file_processing_details": file_stats
        }
        
    except Exception as e:
        # Clean up temporary files in case of error
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except:
                pass
                
        logger.error(f"GPT-4o Document Intelligence API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def gpt4o_multimodal_service(system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
    """OpenAI GPT-4o LLM service function for multimodal content generation with file support"""
    try:
        # Fixed deployment model
        DEPLOYMENT = 'gpt-4o'
        temp_files = []  # Keep track of temporary files for cleanup
        
        # Track file processing statistics
        file_stats = {
            "documents_processed": 0,
            "images_processed": 0,
            "text_files_processed": 0,
            "pages_processed": 0  # Add this field to track page counts
        }
        
        # Track processing errors to report back to user
        processing_errors = []
        
        # Prepare message content
        message_content = []
        
        # First add the user's input text
        message_content.append({
            "type": "text", 
            "text": user_input
        })
        
        # Process file_ids if any
        if file_ids and isinstance(file_ids, list) and len(file_ids) > 0:
            for file_id in file_ids:
                # Get file details using FileService
                file_info, error = FileService.get_file_url(file_id, user_id)
                
                if error:
                    logger.error(f"Error retrieving file {file_id}: {error}")
                    continue
                
                # Check if we have the necessary file information
                if not file_info or 'file_url' not in file_info or 'content_type' not in file_info:
                    logger.error(f"Incomplete file information for file {file_id}")
                    continue
                
                file_url = file_info['file_url']
                content_type = file_info['content_type']
                file_name = file_info.get('file_name', f"file_{file_id}")
                
                try:
                    # Create a temporary file to store the downloaded content
                    fd, temp_path = tempfile.mkstemp(suffix=f'.{file_name.rsplit(".", 1)[1].lower()}' if '.' in file_name else '')
                    temp_files.append(temp_path)
                    
                    # Download the file from Azure Blob Storage
                    import requests
                    response = requests.get(file_url)
                    response.raise_for_status()
                    
                    with os.fdopen(fd, 'wb') as tmp:
                        tmp.write(response.content)
                    
                    # Check file type and process accordingly
                    if content_type.startswith('image/'):
                        # For images, encode as base64 and add to message
                        file_stats["images_processed"] += 1
                        
                        with open(temp_path, 'rb') as img_file:
                            base64_image = base64.b64encode(img_file.read()).decode('utf-8')
                        
                        # Add as image content
                        message_content.append({
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{content_type};base64,{base64_image}"
                            }
                        })
                        
                    elif is_document_intelligence_supported(file_name):
                        # Process document with intelligent extraction based on user query
                        file_stats["documents_processed"] += 1
                        
                        try:
                            # Process document with intelligent extraction
                            logger.info(f"Processing document: {file_name} (type: {content_type})")
                            extracted_content, page_count = process_with_intelligent_extraction(temp_path, file_name, user_input)
                            file_stats["pages_processed"] += page_count  # Track page count
                            
                            # Add to message content
                            message_content.append({
                                "type": "text",
                                "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                            })
                        except Exception as doc_error:
                            error_msg = f"Error processing document {file_name}: {str(doc_error)}"
                            logger.error(error_msg)
                            processing_errors.append(error_msg)
                            message_content.append({
                                "type": "text",
                                "text": f"\n\nError processing document {file_name}: {str(doc_error)}\n"
                            })
                        
                    elif file_name.endswith('.txt') or content_type == 'text/plain':
                        # Process plain text file
                        file_stats["text_files_processed"] += 1
                        
                        # Extract text
                        extracted_content = process_text_file(temp_path)
                        
                        # Add to message content
                        message_content.append({
                            "type": "text",
                            "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                        })
                        
                    else:
                        # For unsupported file types, just mention they were uploaded
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nA file named '{file_name}' was uploaded, but its content type ({content_type}) is not supported for extraction."
                        })
                        
                except requests.RequestException as e:
                    logger.error(f"Error downloading file {file_id}: {str(e)}")
                    message_content.append({
                        "type": "text",
                        "text": f"\n\nFailed to download file {file_name}: {str(e)}"
                    })
                except Exception as e:
                    logger.error(f"Error processing file {file_id}: {str(e)}")
                    message_content.append({
                        "type": "text",
                        "text": f"\n\nError processing file {file_name}: {str(e)}"
                    })
        
        # Create the chat completion request
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": message_content}
        ]
        
        # Make the API call
        response = openai_client.chat.completions.create(
            model=DEPLOYMENT,
            messages=messages,
            temperature=temperature,
            max_tokens=4000,  # Add reasonable limit
            response_format={"type": "json_object"} if json_output else {"type": "text"}
        )
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
        
        # Total files processed
        total_files_processed = file_stats["documents_processed"] + file_stats["images_processed"] + file_stats["text_files_processed"]
        
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except Exception as e:
                logger.error(f"Error removing temporary file {temp_file}: {str(e)}")
        
        # Add any processing errors to the response
        if processing_errors:
            file_stats["processing_errors"] = processing_errors
        
        return {
            "success": True,
            "result": result,
            "model": DEPLOYMENT,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens,
            "files_processed": total_files_processed,
            "file_processing_details": file_stats
        }
        
    except Exception as e:
        # Clean up temporary files in case of error
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except:
                pass
                
        logger.error(f"GPT-4o Multimodal API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def gpt4o_mini_multimodal_service(system_prompt, user_input, temperature=0.5, json_output=False, file_ids=None, user_id=None):
    """OpenAI GPT-4o-mini LLM service function for multimodal content generation with file support"""
    try:
        # Fixed deployment model
        DEPLOYMENT = 'gpt-4o-mini'
        temp_files = []  # Keep track of temporary files for cleanup
        
        # Track file processing statistics
        file_stats = {
            "documents_processed": 0,
            "images_processed": 0,
            "text_files_processed": 0,
            "pages_processed": 0  # Add page count tracking
        }
        
        # Track processing errors to report back to user
        processing_errors = []
        
        # Prepare message content
        message_content = []
        
        # First add the user's input text
        message_content.append({
            "type": "text", 
            "text": user_input
        })
        
        # Process file_ids if any
        if file_ids and isinstance(file_ids, list) and len(file_ids) > 0:
            for file_id in file_ids:
                # Get file details using FileService
                file_info, error = FileService.get_file_url(file_id, user_id)
                
                if error:
                    logger.error(f"Error retrieving file {file_id}: {error}")
                    continue
                
                # Check if we have the necessary file information
                if not file_info or 'file_url' not in file_info or 'content_type' not in file_info:
                    logger.error(f"Incomplete file information for file {file_id}")
                    continue
                
                file_url = file_info['file_url']
                content_type = file_info['content_type']
                file_name = file_info.get('file_name', f"file_{file_id}")
                
                try:
                    # Create a temporary file to store the downloaded content
                    fd, temp_path = tempfile.mkstemp(suffix=f'.{file_name.rsplit(".", 1)[1].lower()}' if '.' in file_name else '')
                    temp_files.append(temp_path)
                    
                    # Download the file from Azure Blob Storage
                    import requests
                    response = requests.get(file_url)
                    response.raise_for_status()
                    
                    with os.fdopen(fd, 'wb') as tmp:
                        tmp.write(response.content)
                    
                    # Check file type and process accordingly
                    if content_type.startswith('image/'):
                        # For images, encode as base64 and add to message
                        file_stats["images_processed"] += 1
                        
                        with open(temp_path, 'rb') as img_file:
                            base64_image = base64.b64encode(img_file.read()).decode('utf-8')
                        
                        # Add as image content
                        message_content.append({
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{content_type};base64,{base64_image}"
                            }
                        })
                        
                    elif is_document_intelligence_supported(file_name):
                        # Process document with intelligent extraction based on user query
                        file_stats["documents_processed"] += 1
                        
                        try:
                            # For mini model, use more aggressive summarization since context is smaller
                            logger.info(f"Processing document for mini model: {file_name} (type: {content_type})")
                            extracted_content, page_count = process_with_intelligent_extraction(temp_path, file_name, user_input)
                            file_stats["pages_processed"] += page_count  # Track page count
                            
                            # Add to message content
                            message_content.append({
                                "type": "text",
                                "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                            })
                        except Exception as doc_error:
                            error_msg = f"Error processing document {file_name}: {str(doc_error)}"
                            logger.error(error_msg)
                            processing_errors.append(error_msg)
                            message_content.append({
                                "type": "text",
                                "text": f"\n\nError processing document {file_name}: {str(doc_error)}\n"
                            })
                        
                    elif file_name.endswith('.txt') or content_type == 'text/plain':
                        # Process plain text file
                        file_stats["text_files_processed"] += 1
                        
                        # Extract text
                        extracted_content = process_text_file(temp_path)
                        
                        # Add to message content
                        message_content.append({
                            "type": "text",
                            "text": f"\n\n--- Content from {file_name} ---\n{extracted_content}\n--- End of {file_name} ---\n"
                        })
                        
                    else:
                        # For unsupported file types, just mention they were uploaded
                        message_content.append({
                            "type": "text",
                            "text": f"\n\nA file named '{file_name}' was uploaded, but its content type ({content_type}) is not supported for extraction."
                        })
                        
                except requests.RequestException as e:
                    logger.error(f"Error downloading file {file_id}: {str(e)}")
                    message_content.append({
                        "type": "text",
                        "text": f"\n\nFailed to download file {file_name}: {str(e)}"
                    })
                except Exception as e:
                    logger.error(f"Error processing file {file_id}: {str(e)}")
                    message_content.append({
                        "type": "text",
                        "text": f"\n\nError processing file {file_name}: {str(e)}"
                    })
        
        # Create the chat completion request
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": message_content}
        ]
        
        # Make the API call
        response = openai_client.chat.completions.create(
            model=DEPLOYMENT,
            messages=messages,
            temperature=temperature,
            max_tokens=2048,  # Lower token limit for mini
            response_format={"type": "json_object"} if json_output else {"type": "text"}
        )
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        cached_tokens = response.usage.cached_tokens if hasattr(response.usage, 'cached_tokens') else 0
        
        # Total files processed
        total_files_processed = file_stats["documents_processed"] + file_stats["images_processed"] + file_stats["text_files_processed"]
        
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except Exception as e:
                logger.error(f"Error removing temporary file {temp_file}: {str(e)}")
        
        return {
            "success": True,
            "result": result,
            "model": DEPLOYMENT,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens,
            "files_processed": total_files_processed,
            "file_processing_details": file_stats
        }
        
    except Exception as e:
        # Clean up temporary files in case of error
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except:
                pass
                
        logger.error(f"GPT-4o-mini Multimodal API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

def deepseek_r1_service(system_prompt, user_input, temperature=0.5, json_output=False, max_tokens=2048):
    """DeepSeek-R1 LLM service function for chain of thought and deep reasoning"""
    try:
        # Fixed endpoint for DeepSeek R1
        ENDPOINT = 'https://deepseek-r1-aiapi.eastus.models.ai.azure.com'
        
        # Initialize Azure Inference client
        client = ChatCompletionsClient(
            endpoint=ENDPOINT,
            credential=AzureKeyCredential(DEEPSEEK_API_KEY)
        )
        
        # Prepare messages for the model
        messages = []
        
        # Add system message if provided
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        
        # Add user message
        messages.append({"role": "user", "content": user_input})
        
        # Prepare payload
        payload = {
            "messages": messages,
            "max_tokens": max_tokens,
            "temperature": temperature
        }
        
        # Add response format if JSON output is requested
        if json_output:
            payload["response_format"] = {"type": "json_object"}
        
        # Make request to LLM
        response = client.complete(payload)
        
        # Extract response data
        result = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        total_tokens = response.usage.total_tokens
        model_name = response.model if hasattr(response, 'model') else "deepseek-r1-aiapi"
        cached_tokens = 0  # Default to 0 as this model doesn't support cached tokens
        
        return {
            "success": True,
            "result": result,
            "model": model_name,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": total_tokens,
            "cached_tokens": cached_tokens
        }
        
    except Exception as e:
        logger.error(f"DeepSeek-R1 API error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }
